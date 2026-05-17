"""
Jabal — Slide 2 (Thesis & Expectations) renderer.

Layout (per design_spec.md):
  1. Header strip
  2. Section label (INVESTMENT THESIS) + Georgia 17 title ("Executive Summary")
  3. Body card — 4-6 sentence thesis paragraph
  4. Section label (Q2 2026 EARNINGS EXPECTATIONS) + table:
        Metric | Jabal Est. | Consensus | Δ | YoY%
        rows for Revenue / EBITDA / EBITDA margin / Net income / EPS / Dividend
  5. Catalysts + Key Risks — two side-by-side cards (3 bullets each)
  6. Numbered "What to watch on the print"
  7. Footer

Data sources used:
  - canonical_store.valuation_forward            → consensus EPS/PE/etc.
  - canonical_store.valuation_historical         → YoY base
  - get_observations_by_provider(commodities)    → commodity context for thesis text
  - get_observations_by_provider(macro)          → IMF/WB context for thesis text
  - canonical_store.rating_split                 → consensus colour in opening sentence
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Optional

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from src.services.jabal_design_tokens import (
    BLACK, GRAY, MUTED, GOLD, POS, NEG, CARD, WHITE,
    FONT_DISPLAY, FONT_UI,
    SZ_SECTION, SZ_KICKER, SZ_VALUE, SZ_VALUE_LG, SZ_LABEL, SZ_BODY,
    SZ_META, SZ_HEADER, SZ_FOOTER, SZ_BULLET_PILL, SZ_TAB_NUM, SZ_TINY,
    PAGE_W_IN, PAGE_H_IN, MARGIN_L, MARGIN_R, CONTENT_W,
    RULE_THICK_PT, BORDER_THICK_PT, LEFT_ACCENT_W_IN,
    in_, signed_color,
)
from src.services.canonical_store import (
    get_all_fields, get_observations_by_provider,
)
from src.services.render_jabal_snapshot import (
    _text, _hrule, _card, _section_label, _header_strip, _footer,
)


# ── Slide 2 sections ──────────────────────────────────────────

def _section_hero(slide, top: float, label: str, title: str):
    """Section label + Georgia 17 title underneath, used for slide 2/3 heroes."""
    _hrule(slide, MARGIN_L, top, CONTENT_W, color=MUTED)
    _text(slide, MARGIN_L, top + 0.10, CONTENT_W, 0.22, label,
          size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    _text(slide, MARGIN_L, top + 0.32, CONTENT_W, 0.50, title,
          font=FONT_DISPLAY, size=SZ_SECTION, color=BLACK)


def _body_card(slide, left: float, top: float, width: float, height: float,
                body: str):
    """Cream-fill card with gold left-accent and the thesis paragraph inside."""
    _card(slide, left, top, width, height, fill=CARD, border=MUTED,
           left_accent=GOLD)
    # Inset 0.20 from accent edge
    _text(slide, left + 0.20, top + 0.10, width - 0.32, height - 0.20,
          body, size=SZ_BODY, color=BLACK)


def _estimates_table(slide, top: float, rows: list[dict],
                      period_label: str = "ESTIMATE"):
    """5-column estimates table: METRIC | Jabal | YoY | QoQ | CONSENSUS.

    `period_label` is the dynamic header for the Jabal-estimate column
    (e.g. 'Q2 2026E'). Column order matches the institutional reference
    layout: analyst estimate first, deltas in the middle, consensus last.

    Row dict shape: {metric, jabal, yoy, qoq, consensus, is_margin}.
    `is_margin` renders YoY/QoQ in basis points; other rows render in %.
    """
    headers = ["METRIC", period_label.upper(), "YoY", "QoQ", "CONSENSUS"]
    col_w   = [2.20, 1.10, 1.00, 1.00, 1.30]
    row_h   = 0.30
    header_top = top
    # Header
    x = MARGIN_L
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        _text(slide, x, header_top, col_w[i] - 0.05, row_h, h,
              size=SZ_LABEL, color=MUTED, all_caps=True, align=align)
        x += col_w[i]
    _hrule(slide, MARGIN_L, header_top + row_h - 0.02, CONTENT_W,
            color=MUTED)
    # Body rows
    for ri, row in enumerate(rows):
        y = header_top + row_h + ri * row_h
        if ri % 2 == 1:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                in_(MARGIN_L), in_(y - 0.02),
                in_(CONTENT_W), in_(row_h))
            band.fill.solid(); band.fill.fore_color.rgb = CARD
            band.line.fill.background()
        is_margin = bool(row.get("is_margin"))
        x = MARGIN_L
        for i, key in enumerate(["metric", "jabal", "yoy", "qoq", "consensus"]):
            val = row.get(key, "—")
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            color = BLACK
            if key in ("yoy", "qoq") and isinstance(val, (int, float)):
                color = signed_color(val)
                # Margin row: render in basis points. Other rows: percent.
                if is_margin:
                    val = f"{val * 100:+.0f} bps"   # input value is in pp; ×100 → bps
                else:
                    val = f"{val:+.1f}%"
            elif val is None:
                val = "—"
            _text(slide, x, y, col_w[i] - 0.05, row_h, str(val),
                  size=SZ_BODY, color=color, align=align)
            x += col_w[i]


def _two_col_pillared_card(slide, top: float, height: float,
                              left_title: str, left_bullets: list[str],
                              right_title: str, right_bullets: list[str]):
    """Side-by-side cards (Catalysts | Key Risks)."""
    card_w = (CONTENT_W - 0.20) / 2
    # LEFT
    _card(slide, MARGIN_L, top, card_w, height, fill=WHITE, border=MUTED,
           left_accent=POS)
    _text(slide, MARGIN_L + 0.18, top + 0.08, card_w - 0.20, 0.30,
          left_title, size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    bullet_y = top + 0.50
    for b in left_bullets[:3]:
        _bullet_dot(slide, MARGIN_L + 0.20, bullet_y + 0.02, color=POS)
        _text(slide, MARGIN_L + 0.42, bullet_y - 0.02, card_w - 0.50, 0.42, b,
              size=SZ_BODY, color=BLACK)
        bullet_y += 0.46
    # RIGHT
    r_left = MARGIN_L + card_w + 0.20
    _card(slide, r_left, top, card_w, height, fill=WHITE, border=MUTED,
           left_accent=NEG)
    _text(slide, r_left + 0.18, top + 0.08, card_w - 0.20, 0.30,
          right_title, size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    bullet_y = top + 0.50
    for b in right_bullets[:3]:
        _bullet_dot(slide, r_left + 0.20, bullet_y + 0.02, color=NEG)
        _text(slide, r_left + 0.42, bullet_y - 0.02, card_w - 0.50, 0.42, b,
              size=SZ_BODY, color=BLACK)
        bullet_y += 0.46


def _bullet_dot(slide, left, top, *, color=GOLD, size=0.12):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  in_(left), in_(top),
                                  in_(size), in_(size))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return dot


def _numbered_list(slide, top: float, items: list[str]):
    """1/2/3 numerals in gold + body text. Used for 'What to watch on the print'."""
    row_h = 0.27
    for i, body in enumerate(items[:5]):
        y = top + i * row_h
        _text(slide, MARGIN_L, y, 0.25, 0.22, str(i + 1),
              size=SZ_BODY, color=GOLD, bold=True)
        _text(slide, MARGIN_L + 0.30, y, CONTENT_W - 0.30, 0.22, body,
              size=SZ_BODY, color=BLACK)


# ── Public entry point ────────────────────────────────────────

@dataclass
class ThesisData:
    """Slide 2 inputs. Built by the orchestrator from canonical_store +
    light templating of commodities / macro context."""
    exec_summary_body: str
    estimates_rows: list[dict]
    estimates_footnote: str
    estimates_subtitle: str = ""  # Subtitle above the table, e.g. "Jabal estimates vs. consensus  ·  SAR millions unless stated"
    estimates_period_label: str = "ESTIMATE"  # Column header for the Jabal-estimate column, e.g. "Q2 2026E"
    catalysts: list[str] = field(default_factory=list)
    risks: list[str] = field(default_factory=list)
    watch_list: list[str] = field(default_factory=list)
    sources_line: str = ""
    analyst_name: str = "Jabal Research"
    gen_date: str = ""
    total_pages: int = 3
    period_heading: str = "Earnings Expectations"


def render_thesis_slide(prs, data: ThesisData):
    blank = next((L for L in prs.slide_layouts if L.name.lower() == "blank"),
                  prs.slide_layouts[-1])
    slide = prs.slides.add_slide(blank)

    _header_strip(slide, 2, "Thesis & Expectations")

    # Investment Thesis hero + body card
    _section_hero(slide, 1.08, "Investment Thesis", "Executive Summary")
    _body_card(slide, MARGIN_L, 1.96, CONTENT_W, 1.85,
                data.exec_summary_body)

    # Period-aware section label. Falls back to a generic title when the
    # orchestrator didn't supply a quarter (e.g. carry-forward case).
    _section_label(slide, MARGIN_L, 3.96, CONTENT_W, data.period_heading or "Earnings Expectations")
    _text(slide, MARGIN_L, 4.26, CONTENT_W, 0.18,
          data.estimates_subtitle or "Jabal estimates vs. consensus  ·  Local currency unless stated",
          size=Pt(9), color=GRAY)
    _estimates_table(slide, 4.48, data.estimates_rows,
                       period_label=data.estimates_period_label or "ESTIMATE")
    _text(slide, MARGIN_L, 6.88, CONTENT_W, 0.18,
          data.estimates_footnote,
          size=Pt(9), color=GRAY)

    # Catalysts + Risks
    _two_col_pillared_card(
        slide, 7.38, 1.95,
        "Catalysts", data.catalysts,
        "Key Risks", data.risks,
    )

    # Numbered list
    _section_label(slide, MARGIN_L, 9.53, CONTENT_W, "What to Watch on the Print")
    _numbered_list(slide, 9.81, data.watch_list)

    _footer(slide, 2, data.total_pages, data.sources_line,
             data.analyst_name, data.gen_date)
    return slide


# ── Data adapter ──────────────────────────────────────────────

def _template_exec_summary(cv: dict, commodities: dict,
                              macro_obs: dict) -> str:
    """Compose a 4–6 sentence thesis paragraph from canonical data.
    This is a generic template; for a polished deck the analyst rewrites
    it, but the data-driven scaffold means it's never blank."""
    name = "the company"
    sector = industry = "—"
    profile = cv.get("company_profile")
    if profile and isinstance(profile.value, dict):
        name = profile.value.get("name") or name
        sector = profile.value.get("sector") or sector
        industry = profile.value.get("industry") or industry

    rating_val = cv.get("rating_split")
    rating = ""
    n_an = 0
    if rating_val and isinstance(rating_val.value, dict):
        # Prettify the provider enum so the fallback paragraph reads
        # "Strong Buy" rather than "strong_buy" / "STRONG_BUY".
        from src.services.render_jabal_snapshot import _pretty_rating
        rating = _pretty_rating(rating_val.value.get("consensus")) or ""
        n_an   = int(rating_val.value.get("total") or 0)

    target = cv.get("target_price")
    target_text = ""
    if target and isinstance(target.value, dict):
        m = target.value.get("mean")
        if m:
            target_text = f"; consensus target {m:.2f}"

    val_hist = cv.get("valuation_historical")
    pe_text = ""
    if val_hist and isinstance(val_hist.value, dict):
        pe = val_hist.value.get("pe")
        if isinstance(pe, list) and any(pe):
            recent = [p for p in pe if isinstance(p, (int, float))]
            if recent:
                pe_text = f" The shares trade at a P/E around {recent[-1]:.1f}x trailing earnings."

    commodity_text = ""
    industry_commodities = (commodities.get("company_profile") or {}).get("industry_commodities", {}) \
        if commodities else {}
    if industry_commodities:
        bits = []
        for tag, info in industry_commodities.items():
            if not isinstance(info, dict):
                continue
            val = info.get("value")
            yoy = info.get("yoy_pct")
            unit = info.get("unit", "")
            if val is None:
                continue
            yoy_str = f" ({yoy:+.1f}% YoY)" if yoy is not None else ""
            bits.append(f"{tag} at {val:.0f} {unit}{yoy_str}")
        if bits:
            commodity_text = " The macro and commodity backdrop is anchored by " \
                + "; ".join(bits[:2]) + "."

    macro_text = ""
    mp = macro_obs.get("company_profile") if macro_obs else None
    if mp:
        gdp_act = mp.get("gdp_growth_pct")
        gdp_fcst = mp.get("gdp_growth_fcst_next_pct")
        infl_fcst = mp.get("inflation_fcst_next_pct")
        parts = []
        if gdp_act is not None:
            parts.append(f"local-economy GDP growth recently at {gdp_act:.1f}%")
        if gdp_fcst is not None:
            parts.append(f"IMF projecting {gdp_fcst:.1f}% next year")
        if infl_fcst is not None:
            parts.append(f"inflation forecast at {infl_fcst:.1f}%")
        if parts:
            macro_text = " Macro context: " + ", ".join(parts) + "."

    rating_line = ""
    if rating and n_an:
        rating_line = f" Street consensus is {rating} ({n_an} analysts covering){target_text}."

    body = (
        f"Jabal maintains a constructive view into the upcoming print. "
        f"{name} sits in the {sector} sector ({industry})."
        f"{rating_line}"
        f"{pe_text}"
        f"{commodity_text}"
        f"{macro_text}"
        f" The thesis hinges on three factors: demand trajectory through H2, "
        f"feedstock/input-cost discipline, and management's tone on guidance "
        f"during the call. We expect the print itself to be within consensus "
        f"bounds; the read on forward commentary is the swing factor."
    )
    return body


def _fmt_money_b(v):
    """Render a raw money value as 'X.YB' / 'X.YT' / 'X,Y00M'. None on bad input."""
    if not isinstance(v, (int, float)):
        return None
    if abs(v) >= 1e12: return f"{v/1e12:,.2f}T"
    if abs(v) >= 1e9:  return f"{v/1e9:,.1f}B"
    if abs(v) >= 1e6:  return f"{v/1e6:,.0f}M"
    return f"{v:,.0f}"


def _prior_year_same_q(period_label: str) -> str | None:
    """'2024-Q3' -> '2023-Q3'; '2024-Q3 (Mar)' -> '2023-Q3 (Mar)'. None on bad input."""
    import re as _re
    m = _re.match(r"^\s*(\d{4})(\D.*)$", str(period_label or ""))
    if not m:
        return None
    try:
        return f"{int(m.group(1)) - 1}{m.group(2)}"
    except ValueError:
        return None


def _yoy_pct(curr, prev) -> float | None:
    try:
        if curr is None or prev in (None, 0):
            return None
        return (float(curr) / float(prev) - 1.0) * 100
    except (TypeError, ValueError, ZeroDivisionError):
        return None


def _yoy_bps(curr_pct, prev_pct) -> float | None:
    """Margin YoY in basis points (return as percentage for unified rendering)."""
    try:
        if curr_pct is None or prev_pct is None:
            return None
        return float(curr_pct) - float(prev_pct)
    except (TypeError, ValueError):
        return None


def _ms_quarterly_split(ms_q: dict | None) -> tuple[dict, dict, dict]:
    """From `ms_quarterly_forecasts.quarterly`, return (next_est, latest_actual,
    prior_year_actual). Each is a dict {metric_key: value}.

    MS interleaves actuals and forecasts in one period list. We split by
    whether the announcement_date is in the past (actual) or future
    (estimate). Falls back to empty dicts on any shape problem.
    """
    empty = ({}, {}, {})
    if not isinstance(ms_q, dict):
        return empty
    q = ms_q.get("quarterly") or {}
    if not isinstance(q, dict):
        return empty
    periods   = q.get("periods")  or []
    net_sales = q.get("net_sales") or []
    ebitda    = q.get("ebitda")   or []
    nii       = q.get("nii")      or []   # banks (when MS publishes it)
    net_inc   = q.get("net_income") or []
    eps       = q.get("eps")      or []
    ann       = q.get("announcement_dates") or []
    n = len(periods)
    if not n:
        return empty
    # Pad short lists with None so zip aligns.
    def _pad(xs): return list(xs) + [None] * (n - len(xs))
    rows = list(zip(periods, _pad(net_sales), _pad(ebitda), _pad(nii),
                     _pad(net_inc), _pad(eps), _pad(ann)))

    from datetime import datetime as _dt
    today = _dt.now().date()

    def _parse_ms_date(date_str: str):
        """MS publishes announcement dates as MM/DD/YY (e.g. '4/23/26').
        Try each known shape in turn and return a `date` or None."""
        s = str(date_str or "").strip()
        if not s or s == "-":
            return None
        for fmt in ("%Y-%m-%d", "%m/%d/%y", "%m/%d/%Y", "%-m/%-d/%y"):
            try:
                return _dt.strptime(s, fmt).date()
            except (ValueError, TypeError):
                continue
        return None

    def _is_estimate(date_str) -> bool:
        """True if the row is a forward forecast (no past announcement)."""
        d = _parse_ms_date(date_str)
        if d is None:
            # Empty / '-' / unparseable → treat as forward estimate.
            return True
        return d >= today

    actuals = [r for r in rows if not _is_estimate(r[6])]
    estimates = [r for r in rows if _is_estimate(r[6])]

    def _to_dict(r):
        p, rev, eb, n_, ni, e, ad = r
        return {
            "period": p, "revenue": rev, "ebitda": eb,
            "nii": n_, "net_income": ni, "eps": e, "ann": ad,
        }

    next_est = _to_dict(estimates[0]) if estimates else {}
    latest_actual = _to_dict(actuals[-1]) if actuals else {}

    # Helper: parse a period label ("Q2 2026" / "2Q26" / "2026-Q2" / "2026Q2")
    # into (year, quarter). Returns None on no match.
    import re as _r
    def _parse_q(period: str):
        s = str(period or "")
        m = (_r.search(r"(\d{4})\s*Q(\d)|Q(\d)\s*(\d{4})", s, _r.I)
             or _r.search(r"(\d{4})-Q(\d)", s, _r.I)
             or _r.search(r"(\d{4})Q(\d)", s, _r.I))
        if not m:
            return None
        y = int(m.group(1) or m.group(4))
        q = int(m.group(2) or m.group(3))
        return (y, q)

    def _find_in_actuals(target_year, target_q):
        for a in actuals:
            yq = _parse_q(a[0])
            if yq == (target_year, target_q):
                return _to_dict(a)
        return None

    # Prior-year-same-Q for the LATEST actual (used for last-reported YoY).
    prior = {}
    if latest_actual.get("period"):
        yq = _parse_q(latest_actual["period"])
        if yq:
            prior = _find_in_actuals(yq[0] - 1, yq[1]) or {}

    # Prior-year-same-Q for the NEXT estimate (used for forecast-YoY).
    prior_of_next = {}
    if next_est.get("period"):
        yq = _parse_q(next_est["period"])
        if yq:
            prior_of_next = _find_in_actuals(yq[0] - 1, yq[1]) or {}

    # Prior-quarter actual: the quarter immediately preceding the next
    # forecast. Needed for QoQ (e.g. Q2 2026 forecast vs Q1 2026 actual).
    prior_quarter = {}
    if next_est.get("period"):
        yq = _parse_q(next_est["period"])
        if yq:
            ny, nq = yq
            if nq > 1:
                prior_quarter = _find_in_actuals(ny, nq - 1) or {}
            else:
                prior_quarter = _find_in_actuals(ny - 1, 4) or {}

    return next_est, latest_actual, prior, prior_of_next, prior_quarter


def _investing_actuals_yoy(ticker: str) -> dict:
    """Pull NEXT forecast vs prior-year-same-quarter actual from Investing's
    earnings page, returning {revenue, eps} YoY percentages.

    Third-tier fallback for the slide-2 forecast-YoY column when MS doesn't
    publish a clean quarterly forecast block. Forecast row is the first
    earnings_history entry where epsActual is None (the upcoming print);
    prior_actual is the entry with reportYear = forecast_year - 1 and the
    same quarter.

    Falls back to (latest actual vs prior-year actual) only when no
    forecast row exists at all (rare; some thinly-covered names have only
    actuals on the Investing page).
    """
    try:
        from src.providers.probe_investing import _fetch_earnings_page, _slug
    except ImportError:
        return {}
    slug = _slug(ticker)
    if not slug:
        return {}
    state = _fetch_earnings_page(slug)
    if not state:
        return {}
    es = state.get("earningsStore") or {}
    if not isinstance(es, dict):
        return {}
    rows = es.get("earnings") or []
    actuals = [r for r in rows
                if isinstance(r, dict)
                and isinstance(r.get("epsActual"), (int, float))
                and isinstance(r.get("reportYear"), int)
                and isinstance(r.get("reportMonth"), int)]
    # Forecast = first row with epsActual None / epsForecast present.
    forecast = None
    for r in rows or []:
        if not isinstance(r, dict): continue
        if r.get("epsActual") is None and isinstance(r.get("reportYear"), int):
            if (isinstance(r.get("epsForecast"), (int, float))
                or isinstance(r.get("revenueForecast"), (int, float))):
                forecast = r
                break

    # Prefer forecast-vs-prior-year-actual when both exist. Falls back to
    # latest-actual-vs-prior-year-actual when no forecast on the page.
    if forecast:
        ny = forecast["reportYear"]
        nq = ((forecast["reportMonth"] - 1) // 3 + 1)
        prior = next((
            r for r in actuals
            if r["reportYear"] == ny - 1
            and ((r["reportMonth"] - 1) // 3 + 1) == nq
        ), None)
        if not prior:
            return {}
        latest = forecast
        # Forecast row uses different key names — adapt below.
        rev_l = forecast.get("revenueForecast")
        rev_p = prior.get("revenueActual")
        eps_l = forecast.get("epsForecast")
        eps_p = prior.get("epsActual")
        out = {}
        if isinstance(rev_l, (int, float)) and isinstance(rev_p, (int, float)) and rev_p:
            out["revenue"] = (rev_l / rev_p - 1.0) * 100
        if isinstance(eps_l, (int, float)) and isinstance(eps_p, (int, float)) and eps_p:
            out["eps"] = (eps_l / eps_p - 1.0) * 100
        return out

    # No forecast row — fall back to last-actual-vs-prior-year-actual.
    if not actuals:
        return {}
    latest = actuals[0]
    latest_year = latest["reportYear"]
    latest_qm = ((latest["reportMonth"] - 1) // 3 + 1)
    prior = next((
        r for r in actuals[1:]
        if r["reportYear"] == latest_year - 1
        and ((r["reportMonth"] - 1) // 3 + 1) == latest_qm
    ), None)
    if not prior:
        return {}
    out = {}
    rev_l, rev_p = latest.get("revenueActual"), prior.get("revenueActual")
    if isinstance(rev_l, (int, float)) and isinstance(rev_p, (int, float)) and rev_p:
        out["revenue"] = (rev_l / rev_p - 1.0) * 100
    eps_l, eps_p = latest.get("epsActual"), prior.get("epsActual")
    if isinstance(eps_l, (int, float)) and isinstance(eps_p, (int, float)) and eps_p:
        eps_yoy = (eps_l / eps_p - 1.0) * 100
        # Investing rounds bank EPS to 2dp (BKMB Q1 2026 = 0.01, Q1 2025 = 0.01
        # → YoY 0.0% even when Net Income grew 9%). When the rounded values
        # match exactly but the underlying business changed, fall back to
        # the Net-Income YoY as the EPS proxy — true absent share-count change.
        if eps_l == eps_p and abs(eps_l) <= 0.05:
            ni_yoy = None
            try:
                from src.services.store_actuals import latest_actuals  # unused but kept for safety
            except ImportError:
                pass
            # Compute NI YoY from the same Investing earnings page if it
            # exposes revenue (a reasonable proxy for NI YoY when bank EPS
            # is rounded). We don't have NI in earnings rows on Investing
            # — instead, just suppress the EPS YoY rather than report a
            # misleading 0.0%.
            out["eps"] = None  # caller renders '—'
        else:
            out["eps"] = eps_yoy
    return out


def _build_estimates_rows(cv: dict, quarterly: list | None = None,
                            is_bank: bool = False,
                            ms_quarterly_forecasts: dict | None = None,
                            ticker: str = "",
                            currency: str = "") -> tuple[list[dict], str]:
    """Build rows for the slide-2 estimates table.

    Returns (rows, unit_suffix). The caller bakes unit_suffix into the
    table subtitle (e.g. "SAR millions unless stated").

    Row schema:
      Non-bank: Revenue / EBITDA / Net Income / EPS / EBITDA Margin
      Bank:     Operating Income / Net Income / EPS

    Columns:
      • JABAL EST   — '—' in the auto-deck; analyst fills in PPT.
      • YoY         — next-Q consensus vs prior-year-same-Q actual (%).
      • QoQ         — next-Q consensus vs immediately-prior-Q actual (%).
                       Margin rows render YoY/QoQ in bps via `is_margin`.
      • CONSENSUS   — next-Q consensus from Investing/MS.
    """
    val_fwd = cv.get("valuation_forward")
    fwd = val_fwd.value if val_fwd and isinstance(val_fwd.value, dict) else {}

    rev_q_consensus = next((fwd.get(k) for k in ("revenue_next_q",)
                            if isinstance(fwd.get(k), (int, float))), None)
    eps_q_consensus = next((fwd.get(k) for k in ("eps_next_q",)
                            if isinstance(fwd.get(k), (int, float))), None)
    ebitda_q_consensus = None
    ni_q_consensus = None
    nii_q_consensus = None

    # MS /finances/ quarterly forecasts — the canonical_store doesn't carry
    # Revenue/EBITDA/NI consensus for the next quarter, but the MS payload
    # does. When the upstream pipeline passed it through, use it to populate
    # the CONSENSUS column AND to compute YoY (latest actual vs prior-year
    # same quarter from the same MS table). Falls back to canonical_store +
    # payload.quarterly_actuals otherwise.
    yoy_rev = yoy_ebitda = yoy_nii = yoy_ni = yoy_eps = yoy_margin = None
    qoq_rev = qoq_ebitda = qoq_nii = qoq_ni = qoq_eps = qoq_margin = None
    used_ms = False
    if ms_quarterly_forecasts:
        next_est, latest, prior, prior_of_next, prior_quarter = _ms_quarterly_split(ms_quarterly_forecasts)
        # MS publishes unit_scale as a *string* ("million", "billion",
        # "thousand"). Map it to a numeric multiplier so the formatter
        # renders absolute values (e.g. 2614M -> "2.6B").
        _UNIT_MULT = {
            "thousand": 1e3, "thousands": 1e3,
            "million":  1e6, "millions":  1e6, "m": 1e6,
            "billion":  1e9, "billions":  1e9, "b": 1e9,
        }
        raw_scale = ms_quarterly_forecasts.get("unit_scale")
        if isinstance(raw_scale, (int, float)):
            unit_scale = float(raw_scale) or 1.0
        else:
            unit_scale = _UNIT_MULT.get(str(raw_scale or "").strip().lower(), 1.0)
        def _scale(v):
            return v * unit_scale if isinstance(v, (int, float)) else v
        if next_est:
            rev_q_consensus    = rev_q_consensus    or _scale(next_est.get("revenue"))
            ebitda_q_consensus = _scale(next_est.get("ebitda"))
            ni_q_consensus     = _scale(next_est.get("net_income"))
            nii_q_consensus    = _scale(next_est.get("nii"))
            eps_q_consensus    = eps_q_consensus    or next_est.get("eps")
        # ANALYTICAL CONTRACT: YoY in a "Q<n> Earnings Expectations" table
        # is the forecast-vs-prior-year-actual comparison — what the next
        # quarter's consensus implies vs the same quarter last year. Falling
        # back to last-reported YoY would mislabel the column.
        # Use same-source pairing: MS forecast vs MS prior-year actual.
        if next_est and prior_of_next:
            yoy_ebitda = _yoy_pct(next_est.get("ebitda"),     prior_of_next.get("ebitda"))
            yoy_nii    = _yoy_pct(next_est.get("nii"),        prior_of_next.get("nii"))
            yoy_ni     = _yoy_pct(next_est.get("net_income"), prior_of_next.get("net_income"))
            # Revenue + EPS YoY come from the Investing fallback below
            # so the YoY denominator matches the source of the displayed
            # CONSENSUS (Investing wins for those rows via trust ladder).
            curr_rev_ms = next_est.get("revenue") or 0
            prev_rev_ms = prior_of_next.get("revenue") or 0
            curr_eb_ms  = next_est.get("ebitda")
            prev_eb_ms  = prior_of_next.get("ebitda")
            curr_margin = (curr_eb_ms / curr_rev_ms * 100) if curr_eb_ms and curr_rev_ms else None
            prev_margin = (prev_eb_ms / prev_rev_ms * 100) if prev_eb_ms and prev_rev_ms else None
            yoy_margin = _yoy_bps(curr_margin, prev_margin)
            used_ms = True
        # QoQ: next-Q forecast vs immediately-prior-quarter actual.
        if next_est and prior_quarter:
            qoq_rev    = _yoy_pct(next_est.get("revenue"),    prior_quarter.get("revenue"))
            qoq_ebitda = _yoy_pct(next_est.get("ebitda"),     prior_quarter.get("ebitda"))
            qoq_nii    = _yoy_pct(next_est.get("nii"),        prior_quarter.get("nii"))
            qoq_ni     = _yoy_pct(next_est.get("net_income"), prior_quarter.get("net_income"))
            curr_rev_q = next_est.get("revenue") or 0
            prev_rev_q = prior_quarter.get("revenue") or 0
            curr_eb_q  = next_est.get("ebitda")
            prev_eb_q  = prior_quarter.get("ebitda")
            curr_margin_q = (curr_eb_q / curr_rev_q * 100) if curr_eb_q and curr_rev_q else None
            prev_margin_q = (prev_eb_q / prev_rev_q * 100) if prev_eb_q and prev_rev_q else None
            qoq_margin = _yoy_bps(curr_margin_q, prev_margin_q)
        if not (next_est and prior_of_next) and latest and prior:
            # No forecast available for the next quarter (e.g. BKMB has no
            # MS Q2 forecast). Leave YoY blank rather than show a misleading
            # last-reported YoY in a forecast-labeled table.
            used_ms = True

    # Yahoo quarterly_actuals fallback (used when MS forecast block was empty
    # or didn't yield a YoY pair).
    if not used_ms and quarterly:
        # Quarterly list is expected to be FinancialPeriod objects (or dicts);
        # render side accepts both since serialization paths differ.
        def _g(rec, key):
            return getattr(rec, key, None) if not isinstance(rec, dict) else rec.get(key)
        recs_by_period = {(_g(r, "period_label") or "").strip(): r for r in quarterly}
        sorted_periods = sorted(recs_by_period.keys(), reverse=True)
        latest = recs_by_period.get(sorted_periods[0]) if sorted_periods else None
        prior_key = _prior_year_same_q(sorted_periods[0]) if sorted_periods else None
        prior = recs_by_period.get(prior_key) if prior_key else None
        if latest and prior:
            yoy_rev    = _yoy_pct(_g(latest, "revenue"),    _g(prior, "revenue"))
            yoy_ebitda = _yoy_pct(_g(latest, "ebitda"),     _g(prior, "ebitda"))
            yoy_nii    = _yoy_pct(_g(latest, "nii"),        _g(prior, "nii"))
            yoy_ni     = _yoy_pct(_g(latest, "net_income"), _g(prior, "net_income"))
            yoy_eps    = _yoy_pct(_g(latest, "eps"),        _g(prior, "eps"))
            # Margin YoY is computed as a bps delta, but we render it as a
            # percent-point change to keep the column format uniform.
            curr_rev = _g(latest, "revenue") or 0
            prev_rev = _g(prior, "revenue") or 0
            curr_eb  = _g(latest, "ebitda")
            prev_eb  = _g(prior, "ebitda")
            curr_margin = (curr_eb / curr_rev * 100) if curr_eb and curr_rev else None
            prev_margin = (prev_eb / prev_rev * 100) if prev_eb and prev_rev else None
            yoy_margin = _yoy_bps(curr_margin, prev_margin)

    # Investing.com earnings page as last-resort YoY (revenue + eps only).
    # Required for yfinance-blocked tickers where MS forecast block lacks
    # historical actuals or the announcement-date split fails.
    if (yoy_rev is None or yoy_eps is None) and ticker:
        inv = _investing_actuals_yoy(ticker)
        if yoy_rev is None and isinstance(inv.get("revenue"), (int, float)):
            yoy_rev = inv["revenue"]
        if yoy_eps is None and isinstance(inv.get("eps"), (int, float)):
            yoy_eps = inv["eps"]

    # Pick a single magnitude unit for the table — keeps values
    # comparable across rows. Based on the largest absolute value
    # across Revenue / EBITDA / Net Income consensus.
    abs_vals = [v for v in (rev_q_consensus, ebitda_q_consensus, ni_q_consensus, nii_q_consensus)
                  if isinstance(v, (int, float))]
    max_abs = max((abs(v) for v in abs_vals), default=0)
    if   max_abs >= 1e12: unit_div, unit_tag = 1e12, "T"
    elif max_abs >= 1e9:  unit_div, unit_tag = 1e9,  "B"
    elif max_abs >= 1e6:  unit_div, unit_tag = 1e6,  "M"
    else:                 unit_div, unit_tag = 1.0,  ""

    cur = (currency or "").upper()
    unit_suffix = (f"{cur}{unit_tag}" if cur and unit_tag else (cur or unit_tag))
    if unit_tag == "T":   subtitle_units = f"{cur} trillions unless stated"
    elif unit_tag == "B": subtitle_units = f"{cur} billions unless stated"
    elif unit_tag == "M": subtitle_units = f"{cur} millions unless stated"
    else:                  subtitle_units = f"{cur} units unless stated".strip()

    def _money_in_unit(v):
        if not isinstance(v, (int, float)): return None
        scaled = v / unit_div
        # Display precision: 1,000-ish numbers as integer; smaller with 1 dp.
        if abs(scaled) >= 100:
            return f"{scaled:,.0f}"
        if abs(scaled) >= 10:
            return f"{scaled:,.1f}"
        return f"{scaled:,.2f}"

    def _eps_fmt(v):
        if not isinstance(v, (int, float)): return None
        return f"{v:,.2f}"

    def _margin_fmt(v):
        if not isinstance(v, (int, float)): return None
        return f"{v:.1f}%"

    def _row(metric: str, jabal_str: str | None, consensus_str: str | None,
              yoy_val, qoq_val, is_margin: bool = False) -> dict:
        return {
            "metric":    metric,
            "jabal":     jabal_str if jabal_str else "—",
            "yoy":       yoy_val if yoy_val is not None else "—",
            "qoq":       qoq_val if qoq_val is not None else "—",
            "consensus": consensus_str if consensus_str else "—",
            "is_margin": is_margin,
        }

    eps_consensus_str = _eps_fmt(eps_q_consensus)

    if is_bank:
        # Bank schema: Operating Income (NII + non-int) / Net Income / EPS.
        # NII / Non-Int / PPOP / Provisions need broker-level data that
        # MS and Investing don't separately publish — Bloomberg upload is
        # the right path for those rows when wanted.
        rows = [
            _row(f"Operating Income ({unit_suffix})",
                  None, _money_in_unit(rev_q_consensus), yoy_rev, qoq_rev),
            _row(f"Net Income ({unit_suffix})",
                  None, _money_in_unit(ni_q_consensus),  yoy_ni,  qoq_ni),
            _row(f"EPS ({cur})",
                  None, eps_consensus_str,               yoy_eps, qoq_eps),
        ]
        return rows, unit_suffix

    # Non-bank: Revenue / EBITDA / Net Income / EPS / EBITDA Margin.
    # EBITDA Margin = EBITDA / Revenue × 100, computed from forecast pair.
    margin_consensus_str = None
    if (isinstance(ebitda_q_consensus, (int, float))
        and isinstance(rev_q_consensus, (int, float))
        and rev_q_consensus > 0):
        margin_consensus_str = _margin_fmt(ebitda_q_consensus / rev_q_consensus * 100)

    rows = [
        _row(f"Revenue ({unit_suffix})",     None, _money_in_unit(rev_q_consensus),    yoy_rev,    qoq_rev),
        _row(f"EBITDA ({unit_suffix})",      None, _money_in_unit(ebitda_q_consensus), yoy_ebitda, qoq_ebitda),
        _row(f"Net Income ({unit_suffix})",  None, _money_in_unit(ni_q_consensus),     yoy_ni,     qoq_ni),
        _row(f"EPS ({cur})",                  None, eps_consensus_str,                  yoy_eps,    qoq_eps),
        _row("EBITDA Margin",                 None, margin_consensus_str,               yoy_margin, qoq_margin,
              is_margin=True),
    ]
    return rows, unit_suffix


def build_thesis_data(ticker: str, *, analyst_name: str = "Jabal Research",
                        gen_date: str = "",
                        catalysts: Optional[list[str]] = None,
                        risks: Optional[list[str]] = None,
                        watch_list: Optional[list[str]] = None,
                        quarterly: Optional[list] = None,
                        is_bank: bool = False,
                        ms_quarterly_forecasts: Optional[dict] = None,
                        period_heading: Optional[str] = None,
                        ) -> ThesisData:
    cv = get_all_fields(ticker)
    commodities_obs = get_observations_by_provider(ticker, "commodities")
    macro_obs       = get_observations_by_provider(ticker, "macro")
    investing_obs   = get_observations_by_provider(ticker, "investing")

    # Prefer the Gemini-generated summary; fall back to the auto-template
    # when the LLM is unavailable (no API key / rate limit / parse failure).
    llm = None
    try:
        from src.services.llm_summary import generate_summary
        llm = generate_summary(ticker)
    except Exception:
        llm = None
    summary = (llm or {}).get("thesis_paragraph") or _template_exec_summary(
        cv, commodities_obs, macro_obs)
    # Look up listing currency from company_master so the table can label
    # values "Revenue (SARM)" / "(AEDM)" etc. Falls back to canonical
    # profile currency when DB lookup misses.
    deck_currency = ""
    try:
        from src.storage.db import load_company as _lc
        cm = _lc(ticker) or {}
        deck_currency = (cm.get("currency") or "").strip()
    except Exception:
        pass
    if not deck_currency:
        prof = cv.get("company_profile")
        if prof and isinstance(prof.value, dict):
            deck_currency = (prof.value.get("currency") or "").strip()
    rows, unit_suffix = _build_estimates_rows(cv, quarterly=quarterly, is_bank=is_bank,
                                                ms_quarterly_forecasts=ms_quarterly_forecasts,
                                                ticker=ticker, currency=deck_currency)

    # Compose the table footnote: lead with the next-Q anchor (date,
    # consensus source, analyst count) since that's the strongest data
    # point we have for the full panel.
    val_fwd = cv.get("valuation_forward")
    rs = cv.get("rating_split")
    fwd_dict = val_fwd.value if val_fwd and isinstance(val_fwd.value, dict) else {}
    rs_dict  = rs.value if rs and isinstance(rs.value, dict) else {}
    n_an = int(rs_dict.get("total", 0) or 0)
    nq_period = fwd_dict.get("next_q_period") or ""
    nq_date   = fwd_dict.get("next_q_report_date") or ""
    fwd_source = (val_fwd.canonical_source if val_fwd else "—").title()
    footnote_bits = []
    if nq_period and nq_date:
        footnote_bits.append(f"Next print: {nq_date} (period {nq_period})")
    footnote_bits.append(f"Consensus: {fwd_source}")
    if n_an:
        footnote_bits.append(f"{n_an} analysts covering")
    estimates_footnote = "  ·  ".join(footnote_bits)

    # Surface Investing's surprise history as a track-record catalyst line.
    surprise = (investing_obs.get("income_statement_quarterly") or {}).get(
        "surprise_history", [])
    track_record_catalyst = None
    if surprise:
        beats = sum(1 for r in surprise[:4]
                     if isinstance(r.get("eps_surprise_pct"), (int, float))
                     and r["eps_surprise_pct"] > 0)
        n = min(4, len(surprise))
        last = surprise[0]
        last_dir = "beat" if (last.get("eps_surprise_pct") or 0) > 0 else "missed"
        last_pct = abs(last.get("eps_surprise_pct") or 0)
        track_record_catalyst = (
            f"EPS {last_dir} consensus by {last_pct:.1f}% last quarter; "
            f"{beats} of last {n} quarters above estimates"
        )
    default_catalysts = [
        track_record_catalyst or
        "Positive quarterly surprise consistent with prior track record",
        "Constructive guidance on H2 demand and pricing outlook",
        "Capex/project milestones tracking on schedule",
    ]

    # LLM output, when present, replaces every default bullet list too —
    # otherwise the deck mixes a fresh LLM thesis with stale boilerplate.
    llm_catalysts = (llm or {}).get("catalysts") or []
    llm_risks     = (llm or {}).get("risks") or []
    llm_watch     = (llm or {}).get("watch_list") or []

    # Derive the Jabal-estimate column header (e.g. "Q2 2026E") from the
    # period_heading parameter ("Q2 2026 Earnings Expectations" → "Q2 2026E").
    estimates_period_label = "ESTIMATE"
    if period_heading:
        import re as _re_ph
        m = _re_ph.match(r"(Q\d\s+\d{4})\b", period_heading)
        if m:
            estimates_period_label = f"{m.group(1)}E"
        else:
            estimates_period_label = period_heading.upper()
    # Subtitle line under the section heading.
    subtitle_unit_phrase = (
        f"{(deck_currency or '').upper()} {('trillions' if unit_suffix.endswith('T') else 'billions' if unit_suffix.endswith('B') else 'millions' if unit_suffix.endswith('M') else 'units')} unless stated".strip()
    )
    estimates_subtitle = f"Jabal estimates vs. consensus  ·  {subtitle_unit_phrase}"

    # Footnote: source + analyst count + "Bps = basis points" disclosure when
    # margin row is present.
    consensus_source = fwd_source if fwd_source and fwd_source != "—" else "MarketScreener"
    footnote_bits = [f"Estimates: Jabal Research", f"Consensus: {consensus_source}"]
    if n_an:
        footnote_bits[-1] = f"Consensus: {consensus_source} ({n_an} analysts)"
    if not is_bank:
        footnote_bits.append("Bps = basis points")
    estimates_footnote = "  ·  ".join(footnote_bits)

    from datetime import datetime
    return ThesisData(
        exec_summary_body=summary,
        estimates_rows=rows,
        estimates_footnote=estimates_footnote,
        estimates_subtitle=estimates_subtitle,
        estimates_period_label=estimates_period_label,
        catalysts=catalysts or llm_catalysts or default_catalysts,
        risks=risks or llm_risks or [
            "Cautious management tone could validate target-price gap",
            "Feedstock / input cost volatility pressures margin trajectory",
            "Macro / commodity-price softness weighs on top-line growth",
        ],
        watch_list=watch_list or llm_watch or [
            "Forward demand commentary — Q3 order book and pricing trajectory",
            "Feedstock cost outlook and supply-chain commentary",
            "Updated capex schedule and any project-pipeline updates",
        ],
        sources_line=_sources_line_from_cv(cv),
        analyst_name=analyst_name,
        gen_date=gen_date or datetime.utcnow().strftime("%d %b %Y"),
        period_heading=(period_heading or "Earnings Expectations"),
    )


def _sources_line_from_cv(cv: dict) -> str:
    """Re-export of render_jabal_snapshot._sources_line so all three
    slide builders attribute the same way (every contributing provider,
    not just per-field winners)."""
    from src.services.render_jabal_snapshot import _sources_line
    return _sources_line(cv)
