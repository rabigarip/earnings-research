"""
Jabal — Slide 3 (Valuation & Positioning) renderer.

Layout:
  1. Header strip
  2. Section hero (MARKET POSITIONING / "Valuation & Market View")
  3. Two-up chart row:
       LEFT  — 52-week price chart (line)
       RIGHT — P/E multiple 5-year range (horizontal bars)
  4. Peer comparables table
  5. Sentiment row (3 cards): Consensus distribution | Avg target | Last 3 broker actions
  6. Footer

Charts are drawn natively in pptx (line + horizontal bars) rather than
embedded matplotlib images — keeps the file lean and editable.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from datetime import datetime
from typing import Optional

log = logging.getLogger(__name__)

from pptx.enum.shapes import MSO_SHAPE

from src.services.render_jabal_snapshot import _sources_line
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from src.services.jabal_design_tokens import (
    BLACK, GRAY, MUTED, GOLD, GOLD_DK, POS, NEG, CARD, WHITE,
    FONT_DISPLAY, FONT_UI,
    SZ_SECTION, SZ_KICKER, SZ_VALUE, SZ_LABEL, SZ_BODY,
    SZ_HEADER, SZ_FOOTER, SZ_BULLET_PILL, SZ_TINY,
    PAGE_W_IN, PAGE_H_IN, MARGIN_L, MARGIN_R, CONTENT_W,
    RULE_THICK_PT, BORDER_THICK_PT, LEFT_ACCENT_W_IN,
    in_, signed_color,
)
from src.services.canonical_store import (
    get_all_fields, get_observations_by_provider,
)
from src.services.render_jabal_snapshot import (
    _text, _hrule, _card, _section_label, _header_strip, _footer,
)
from src.services.render_jabal_thesis import _section_hero


# ── Native pptx charts ────────────────────────────────────────

def _line_chart_52w(slide, left: float, top: float, width: float, height: float,
                     close_series: list[dict], currency: str = ""):
    """Draw a simple line chart using add_chart with XL_CHART_TYPE.LINE.
    close_series is a list of {date, close} dicts (sparse OK)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    if not close_series:
        _text(slide, left, top + height * 0.45, width, 0.30,
              "No price history available", size=SZ_BODY, color=MUTED,
              align=PP_ALIGN.CENTER)
        return
    cats = [pt["date"][-5:] for pt in close_series]
    vals = [pt["close"] for pt in close_series]
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Close", vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, in_(left), in_(top), in_(width), in_(height), cd
    )
    chart = chart_shape.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    # Style: single gold series, thin line
    try:
        ser = plot.series[0]
        from pptx.dml.color import RGBColor
        line = ser.format.line
        line.color.rgb = GOLD
        line.width = Pt(1.5)
        # No markers — clean line
        from pptx.oxml.ns import qn
        sp_pr = ser._element.find(qn("c:spPr"))
        if sp_pr is None:
            pass
    except (AttributeError, KeyError, TypeError, ValueError):
        pass
    # Hide gridlines for tightness
    try:
        for axis in (chart.value_axis, chart.category_axis):
            axis.major_unit = None
            axis.minor_unit = None
            axis.format.line.fill.background()
            for tl in (axis.tick_labels,):
                tl.font.size = Pt(7)
                tl.font.name = FONT_UI
                from pptx.dml.color import RGBColor
                tl.font.color.rgb = MUTED
    except (AttributeError, KeyError, TypeError, ValueError):
        pass


def _pe_range_chart(slide, left: float, top: float, width: float, height: float,
                     periods: list[str], pe_vals: list[Optional[float]],
                     current_pe: Optional[float],
                     *, draw_header: bool = True):
    """Horizontal bars representing P/E across FY periods, plus a
    'current' diamond marker on the most-recent bar.

    Renders manually with shapes — no chart object — so we can pixel-tune
    the look to match the spec deck.

    `draw_header` is True for legacy stand-alone callers; the new
    `render_valuation_slide` already paints its own header above the
    chart and passes False to avoid the duplicate label.
    """
    if draw_header:
        _text(slide, left, top, width, 0.22, "P/E MULTIPLE  ·  5-YEAR RANGE",
              size=SZ_LABEL, color=MUTED, all_caps=True, bold=True)
    rows = [(p, v) for p, v in zip(periods, pe_vals) if isinstance(v, (int, float))]
    if not rows:
        _text(slide, left, top + height * 0.45, width, 0.30,
              "No P/E history available", size=SZ_BODY, color=MUTED,
              align=PP_ALIGN.CENTER)
        return

    # Axis scale: round nice limits around the data
    vals = [v for _, v in rows]
    lo = max(0, min(vals) * 0.8)
    hi = max(vals) * 1.15
    if hi <= lo:
        hi = lo + 1
    bar_area_left = left + 0.55
    bar_area_w    = width - 0.65
    row_h = 0.30
    # When the caller drew its own header above us, top already points
    # at the start of the chart area — don't push the rows down by an
    # additional 0.30 (which would assume our own header sits at `top`).
    row_top = top + (0.30 if draw_header else 0.0)
    for i, (period, pe) in enumerate(rows):
        y = row_top + i * row_h
        _text(slide, left, y, 0.55, 0.20, period,
              size=SZ_BODY, color=BLACK)
        # Bar: from lo to pe (normalised within axis lo..hi)
        frac_lo = 0
        frac_hi = max(0.02, (pe - lo) / (hi - lo))
        bar_w = bar_area_w * frac_hi
        bar_left = bar_area_left
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            in_(bar_left), in_(y + 0.06),
            in_(bar_w), in_(0.08))
        bar.fill.solid(); bar.fill.fore_color.rgb = CARD
        bar.line.fill.background()
        # Endpoint dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            in_(bar_left + bar_w - 0.05), in_(y + 0.04),
            in_(0.10), in_(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = GOLD_DK
        dot.line.fill.background()
        # P/E label to the right
        _text(slide, bar_left + bar_w + 0.05, y + 0.02, 0.45, 0.20,
              f"{pe:.1f}x", size=SZ_BODY, color=BLACK)

    # Axis ticks below
    axis_y = row_top + len(rows) * row_h + 0.04
    n_ticks = 4
    for i in range(n_ticks):
        frac = i / (n_ticks - 1)
        tick_x = bar_area_left + bar_area_w * frac - 0.15
        tick_val = lo + (hi - lo) * frac
        _text(slide, tick_x, axis_y, 0.40, 0.18, f"{tick_val:.0f}x",
              size=SZ_TINY, color=MUTED)

    # Current marker + legend
    if current_pe is not None and hi > lo:
        cur_frac = max(0.0, min(1.0, (current_pe - lo) / (hi - lo)))
        cur_x = bar_area_left + bar_area_w * cur_frac
        marker = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
            in_(cur_x - 0.08), in_(row_top + len(rows) * row_h - 0.02),
            in_(0.16), in_(0.16))
        marker.fill.solid(); marker.fill.fore_color.rgb = BLACK
        marker.line.fill.background()
        _text(slide, cur_x + 0.10, row_top + len(rows) * row_h - 0.04,
               1.4, 0.18,
               f"Current  ({current_pe:.1f}x)",
               size=SZ_TINY, color=BLACK)


# ── Peer table ────────────────────────────────────────────────

def _peer_label(ticker: str, sector: str = "", industry: str = "") -> str:
    """Derive a peer-table sub-label from the ticker's exchange suffix and
    its sector/industry. Returns a phrase like 'GCC bank peers',
    'India banking peers', 'China internet peers'. Falls back to
    'Selected global peers'."""
    suf = ticker.rsplit(".", 1)[-1].upper() if "." in ticker else ""
    s = (sector or "").lower()
    i = (industry or "").lower()
    text = f"{s} {i}"
    is_bank = any(k in text for k in ("bank", "financial services", "diversified financial"))
    is_oil  = any(k in text for k in ("oil", "gas", "energy", "petroleum", "refining"))
    is_chem = "chem" in text or "fertili" in text or "materials" in text or "basic materials" in text
    is_telecom = "telecom" in text or "communication services" in text
    is_internet = any(k in text for k in ("internet", "interactive media", "software—internet"))
    is_auto = "auto" in text or "vehicle" in text
    is_metal = "steel" in text or "metal" in text or "mining" in text
    gcc = {"SR","AE","QA","OM","KW","BH"}
    india = {"NS","BO"}
    china = {"HK","SS","SZ"}
    if suf in gcc:
        region = "GCC"
    elif suf in india:
        region = "India"
    elif suf in china:
        region = "China/HK" if suf == "HK" else "China"
    else:
        region = "Global"
    if is_bank:
        return f"{region} bank peers"
    if is_oil:
        return "Global oil & gas peers" if region == "Global" else f"{region} oil & gas peers"
    if is_chem:
        return "Global chemicals peers" if region == "Global" else f"{region} chemicals peers"
    if is_internet:
        return "China/HK internet peers" if region in ("China/HK", "China") else f"{region} internet peers"
    if is_auto:
        return f"{region} auto peers"
    if is_metal:
        return f"{region} metals & mining peers"
    if is_telecom:
        return f"{region} telecom peers"
    return f"Selected {region.lower()} peers" if region != "Global" else "Selected global peers"


def _peer_table(slide, top: float, peers: list[dict], *, is_bank: bool = False):
    """Rows: name, ticker, mcap, P/E, then a sector-appropriate book/
    enterprise multiple set, dividend yield, 1Y return.

    Bank schema (7 cols): COMPANY · TICKER · MCAP · P/E · P/TBV · DIV YIELD · 1Y RETURN
      P/TBV is Yahoo's priceToBook (close enough for relative bank
      comps; yfinance doesn't expose tangible book separately). Showing
      a redundant P/B alongside would just consume slide width.

    Non-bank schema (8 cols): COMPANY · TICKER · MCAP · P/E · P/B · EV/EBITDA · DIV YIELD · 1Y RETURN

    Compact, borderless, alternating row tint.
    """
    if is_bank:
        headers = ["COMPANY", "TICKER", "MCAP", "P/E", "P/TBV", "DIV YIELD", "1Y RETURN"]
        col_w   = [1.85, 1.00, 1.05, 0.65, 0.70, 0.85, 0.85]
    else:
        headers = ["COMPANY", "TICKER", "MCAP", "P/E", "P/B", "EV/EBITDA", "DIV YIELD", "1Y RETURN"]
        col_w   = [1.65, 0.95, 0.95, 0.55, 0.55, 0.75, 0.70, 0.70]
    row_h   = 0.28
    # Header
    x = MARGIN_L
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i < 2 else PP_ALIGN.RIGHT
        _text(slide, x, top, col_w[i] - 0.05, row_h, h,
              size=SZ_LABEL, color=MUTED, all_caps=True, align=align)
        x += col_w[i]
    _hrule(slide, MARGIN_L, top + row_h - 0.02, CONTENT_W, color=MUTED)

    # 1Y RETURN is always the last column — track its index for the
    # sign-aware colour and the alignment loop below.
    ret_col_idx = len(headers) - 1

    for ri, p in enumerate(peers[:5]):
        y = top + row_h + ri * row_h
        if ri % 2 == 1:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                in_(MARGIN_L), in_(y - 0.02),
                in_(CONTENT_W), in_(row_h))
            band.fill.solid(); band.fill.fore_color.rgb = CARD
            band.line.fill.background()
        x = MARGIN_L
        if is_bank:
            cells = [
                p.get("name", "—"),
                p.get("ticker", "—"),
                p.get("market_cap_fmt", "—"),
                p.get("pe_fmt", "—"),
                p.get("pb_fmt", "—"),   # rendered under the P/TBV header
                p.get("div_yield_fmt", "—"),
                p.get("ret_1y_fmt", "—"),
            ]
        else:
            cells = [
                p.get("name", "—"),
                p.get("ticker", "—"),
                p.get("market_cap_fmt", "—"),
                p.get("pe_fmt", "—"),
                p.get("pb_fmt", "—"),
                p.get("ev_ebitda_fmt", "—"),
                p.get("div_yield_fmt", "—"),
                p.get("ret_1y_fmt", "—"),
            ]
        ret_val = p.get("ret_1y")
        for i, cell in enumerate(cells):
            align = PP_ALIGN.LEFT if i < 2 else PP_ALIGN.RIGHT
            color = BLACK
            if i == ret_col_idx and isinstance(ret_val, (int, float)):
                color = signed_color(ret_val)
            _text(slide, x, y, col_w[i] - 0.05, row_h, str(cell),
                  size=SZ_BODY, color=color, align=align)
            x += col_w[i]


# ── Sentiment row ─────────────────────────────────────────────

def _sentiment_row(slide, top: float, *, rating_split: dict,
                     n_analysts: int, target_mean: Optional[float],
                     target_range: Optional[tuple], target_implied_pct: Optional[float],
                     broker_actions: list[dict], currency: str):
    card_w = (CONTENT_W - 0.40) / 3
    card_h = 1.15

    # Card 1: Consensus distribution
    _card(slide, MARGIN_L, top, card_w, card_h, fill=WHITE)
    _text(slide, MARGIN_L + 0.12, top + 0.10, card_w - 0.20, 0.18,
          "CONSENSUS DISTRIBUTION", size=SZ_LABEL, color=MUTED,
          all_caps=True, bold=True)
    # Three bars proportional to buy/hold/sell
    total = max(1, sum(rating_split.values())) if rating_split else 1
    seg_top = top + 0.34
    bar_h   = 0.16
    inner_x = MARGIN_L + 0.12
    inner_w = card_w - 0.24
    buy_w  = inner_w * (rating_split.get("buy",  0) / total)
    hold_w = inner_w * (rating_split.get("hold", 0) / total)
    sell_w = inner_w * (rating_split.get("sell", 0) / total)
    for w, color, x_off in [(buy_w, POS, 0),
                              (hold_w, GOLD, buy_w),
                              (sell_w, NEG, buy_w + hold_w)]:
        if w <= 0.001:
            continue
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            in_(inner_x + x_off), in_(seg_top),
            in_(w), in_(bar_h))
        seg.fill.solid(); seg.fill.fore_color.rgb = color
        seg.line.fill.background()
    # Labels below the bar
    _text(slide, inner_x, seg_top + bar_h + 0.06, card_w - 0.24, 0.18,
          f"Buy {int(rating_split.get('buy',0)/total*100)}%   "
          f"Hold {int(rating_split.get('hold',0)/total*100)}%   "
          f"Sell {int(rating_split.get('sell',0)/total*100)}%",
          size=Pt(9), color=GRAY)
    _text(slide, inner_x, seg_top + bar_h + 0.30, card_w - 0.24, 0.20,
          f"{n_analysts} analysts covering",
          size=SZ_LABEL, color=GRAY)

    # Card 2: Average target price
    c2_left = MARGIN_L + card_w + 0.20
    _card(slide, c2_left, top, card_w, card_h, fill=WHITE)
    _text(slide, c2_left + 0.12, top + 0.10, card_w - 0.20, 0.18,
          "AVERAGE TARGET PRICE", size=SZ_LABEL, color=MUTED,
          all_caps=True, bold=True)
    target_str = (f"{currency} {target_mean:,.2f}" if target_mean is not None
                   else "—")
    _text(slide, c2_left + 0.12, top + 0.32, card_w - 0.20, 0.40,
          target_str, size=Pt(16), color=BLACK, bold=True)
    if target_range:
        # Pick decimal precision by absolute value so a low-priced OMR
        # ticker doesn't round to "OMR 0 — 1". Sub-1 → 3dp; 1-10 → 2dp;
        # 10-100 → 1dp; >=100 → 0dp.
        def _rng_fmt(v):
            if v is None: return ""
            v_abs = abs(float(v))
            if v_abs < 1:   return f"{v:,.3f}"
            if v_abs < 10:  return f"{v:,.2f}"
            if v_abs < 100: return f"{v:,.1f}"
            return f"{v:,.0f}"
        rng_str = (f"Range  {currency} {_rng_fmt(target_range[0])} — {_rng_fmt(target_range[1])}"
                    if target_range[0] is not None and target_range[1] is not None
                    else "")
        _text(slide, c2_left + 0.12, top + 0.74, card_w - 0.20, 0.18,
              rng_str, size=Pt(9), color=GRAY)
    if target_implied_pct is not None:
        _text(slide, c2_left + 0.12, top + 0.92, card_w - 0.20, 0.18,
              f"Implied {target_implied_pct:+.1f}% vs last close",
              size=Pt(9), color=signed_color(target_implied_pct))

    # Card 3: Last broker actions
    c3_left = MARGIN_L + 2 * (card_w + 0.20)
    _card(slide, c3_left, top, card_w, card_h, fill=WHITE)
    _text(slide, c3_left + 0.12, top + 0.10, card_w - 0.20, 0.18,
          "LAST BROKER ACTIONS", size=SZ_LABEL, color=MUTED,
          all_caps=True, bold=True)
    if not broker_actions:
        _text(slide, c3_left + 0.12, top + 0.40, card_w - 0.20, 0.20,
              "No broker actions in feed",
              size=Pt(9), color=MUTED)
    else:
        for i, ba in enumerate(broker_actions[:3]):
            y = top + 0.34 + i * 0.26
            _text(slide, c3_left + 0.12, y, 0.55, 0.20,
                  ba.get("date", "—"), size=SZ_BODY, color=GRAY)
            _text(slide, c3_left + 0.70, y, card_w - 0.78, 0.20,
                  ba.get("text", "—"), size=SZ_BODY, color=BLACK)


# ── Public entry point ────────────────────────────────────────

@dataclass
class ValuationData:
    company_name: str
    close_series: list[dict]
    currency: str
    pe_periods: list[str]
    pe_values: list[Optional[float]]
    pe_current: Optional[float]
    peers: list[dict]
    peer_table_label: str
    rating_split: dict
    n_analysts: int
    target_mean: Optional[float]
    target_range: Optional[tuple]
    target_implied_pct: Optional[float]
    broker_actions: list[dict]
    sources_line: str
    analyst_name: str
    gen_date: str
    total_pages: int = 3
    is_bank: bool = False
    # Earnings-history chart inputs (replaces the legacy "Market Sentiment"
    # row at the bottom of slide 3 — duplicates info already on slide 1).
    surprise_history: list[dict] | None = None
    # "EPS" / "Net Income" / "Net Sales" — depends on which row MS's
    # /calendar/ quarterly_results carried for the ticker. Used by the
    # chart title so a bank-with-no-EPS-row deck doesn't claim it's
    # plotting EPS surprises when it's actually net-income surprises.
    surprise_metric_label: str = "EPS"
    ticker: str = ""
    # Optional 5y forward-P/E history (for the new historical-range view).
    # When empty, the P/E chart falls back to the forecast-bar view.
    pe_history: list[dict] | None = None


def render_valuation_slide(prs, data: ValuationData):
    blank = next((L for L in prs.slide_layouts if L.name.lower() == "blank"),
                  prs.slide_layouts[-1])
    slide = prs.slides.add_slide(blank)

    _header_strip(slide, 3, "Valuation & Positioning")
    _section_hero(slide, 1.08, "Market Positioning",
                    "Valuation & Market View")

    # Two-up chart row — matplotlib PNGs embedded as pictures so we get
    # cleaner date axes, range shading, and current-multiple markers than
    # the native python-pptx chart objects allow.
    from io import BytesIO
    from src.services.render_charts_mpl import (
        render_52w_price_chart, render_pe_historical_chart,
        render_earnings_history_chart,
    )

    chart_top = 1.96
    chart_h = 2.40
    col_w = (CONTENT_W - 0.20) / 2

    # Left: 52w price chart
    _text(slide, MARGIN_L, chart_top, col_w, 0.22,
          f"52-WEEK PRICE  ·  {data.currency}",
          size=SZ_LABEL, color=MUTED, all_caps=True, bold=True)
    price_png = render_52w_price_chart(data.close_series, currency=data.currency)
    if price_png:
        slide.shapes.add_picture(BytesIO(price_png),
                                   in_(MARGIN_L), in_(chart_top + 0.26),
                                   width=in_(col_w), height=in_(chart_h - 0.30))
    else:
        # Matplotlib unavailable or empty series — keep the native fallback
        # so the slide still renders.
        _line_chart_52w(slide, MARGIN_L, chart_top + 0.26, col_w, chart_h - 0.30,
                         data.close_series, currency=data.currency)

    # Right: forward P/E across FY periods (Bloomberg "Current Multiples"
    # convention — today's price ÷ each FY's EPS estimate, projected
    # forward). Falls back to native bar-chart when matplotlib returns
    # None, but we DON'T also re-render the header in that branch
    # (`_pe_range_chart` historically painted its own header, which gave
    # us a double-label on tickers with empty P/E data).
    right_left = MARGIN_L + col_w + 0.20
    _text(slide, right_left, chart_top, col_w, 0.22,
          "FORWARD P/E  ·  CURRENT MULTIPLES",
          size=SZ_LABEL, color=MUTED, all_caps=True, bold=True)
    pe_png = render_pe_historical_chart(
        data.pe_history, data.pe_periods, data.pe_values, data.pe_current,
    )
    if pe_png:
        slide.shapes.add_picture(BytesIO(pe_png),
                                   in_(right_left), in_(chart_top + 0.26),
                                   width=in_(col_w), height=in_(chart_h - 0.30))
    else:
        # Native fallback — but suppress its internal header to avoid
        # the duplicate-label bug observed on BKMB.
        _pe_range_chart(slide, right_left, chart_top + 0.26, col_w, chart_h - 0.30,
                         data.pe_periods, data.pe_values, data.pe_current,
                         draw_header=False)

    # Peer table
    _section_label(slide, MARGIN_L, 4.97, CONTENT_W,
                    f"Peer Comparables  ·  {data.peer_table_label}")
    _peer_table(slide, 5.29, data.peers, is_bank=data.is_bank)

    # Earnings history (replaces the legacy "Market Sentiment" row — the
    # consensus distribution + average target are already on slide 1, and
    # the broker-actions card was almost always empty for our universe).
    _eh_metric = data.surprise_metric_label or "EPS"
    _section_label(slide, MARGIN_L, 8.35, CONTENT_W,
                    f"Earnings History  ·  {_eh_metric} Actual vs Estimate")
    eh_png = render_earnings_history_chart(
        data.surprise_history or [],
        price_series=data.close_series,
        ticker=data.ticker,
        currency=data.currency,
        max_quarters=8,
        metric_label=_eh_metric,
    )
    if eh_png:
        slide.shapes.add_picture(BytesIO(eh_png),
                                   in_(MARGIN_L), in_(8.65),
                                   width=in_(CONTENT_W), height=in_(2.50))
    else:
        _text(slide, MARGIN_L, 8.85, CONTENT_W, 0.30,
              "No earnings surprise history available",
              size=SZ_BODY, color=MUTED, align=PP_ALIGN.CENTER)

    _footer(slide, 3, data.total_pages, data.sources_line,
             data.analyst_name, data.gen_date)
    return slide


# ── Data adapter ──────────────────────────────────────────────

def _load_ms_calendar_from_snapshot(ticker: str) -> dict | None:
    """Load `quarterly_results` directly from the committed MS calendar
    snapshot under data/marketscreener/, bypassing the pipeline. Used
    when `payload.ms_calendar_events` arrives as None despite the
    snapshot existing on disk (common when fetch_marketscreener_pages
    flags PARTIAL because live HTTP 403'd, even though the snapshot
    fallback fed the parser successfully).

    Resolves the ticker → MS slug → ISIN via company_master, builds the
    runtime cache prefix (`ms_<ticker>_<isin>_<slug>`), points at the
    snapshot file, parses with the production parser, and returns the
    `quarterly_results` block.

    Returns None when company isn't curated for MS, the snapshot file
    is absent, or parsing fails.
    """
    try:
        from src.storage.db import load_company
    except Exception:
        return None
    row = load_company(ticker) or {}
    slug = (row.get("marketscreener_id") or "").strip()
    isin = (row.get("isin") or "").strip() or "noisin"
    if not slug:
        return None
    # Match the cache slug shape used by the production fetcher path.
    t_safe = ticker.replace(".", "_").strip() or "unknown"
    cache_slug = f"ms_{t_safe}_{isin}_{slug}_calendar"
    import re as _re_load
    safe = _re_load.sub(r"[^a-zA-Z0-9-]", "_", cache_slug)[:80]
    try:
        from src.config import root
        snapshot_path = root() / "data" / "marketscreener" / f"ms_{safe}.html"
    except Exception:
        return None
    if not snapshot_path.exists():
        return None
    try:
        html = snapshot_path.read_text(encoding="utf-8")
    except Exception:
        return None
    try:
        from bs4 import BeautifulSoup
        from src.providers.marketscreener_pages import _parse_quarterly_results_table
        soup = BeautifulSoup(html, "lxml")
        qr = _parse_quarterly_results_table(soup)
    except Exception:
        return None
    if not (qr.get("quarters") and qr.get("rows")):
        return None
    return qr


def _surprise_history_from_ms_calendar(
        ms_calendar_events: dict | None) -> tuple[list[dict], str]:
    """Derive a surprise history from MS's /calendar/ quarterly_results.

    Investing.com is the primary source for `surprise_history`, but for
    thinly-covered names like BKMB.OM Investing has the announce dates
    only — no actuals / estimates. MS publishes released + forecast
    per metric per quarter on the /calendar/ page.

    Picks the BEST available metric in priority order:
      1. EPS         — preferred, matches Investing's chart convention
      2. Net Income  — banks (BKMB, ENBD, etc.) — MS calendar has no EPS row
      3. Net Sales   — fallback for tickers where even NI isn't published
                       (the user's MS screenshot for BKMB literally shows
                       "Quarterly REVENUE — Rate of surprise")

    Returns (rows, metric_label) where metric_label is one of "EPS",
    "Net Income", "Net Sales" so the chart can adapt its title /
    y-axis. The data shape keeps the historical "eps_actual" /
    "eps_estimate" / "eps_surprise_pct" key names so render_charts_mpl
    consumes it unchanged — the prefix is a relic, the underlying
    semantic is "actual vs estimate for whatever metric we found".

    Returns ([], "") when MS calendar data is unavailable.
    """
    if not isinstance(ms_calendar_events, dict):
        return [], ""
    qr = ms_calendar_events.get("quarterly_results") or {}
    quarters = qr.get("quarters") or []
    rows = qr.get("rows") or []
    if not (quarters and rows):
        return [], ""

    # Priority cascade: EPS → Net Income → Net Sales.
    def _find(*keys):
        norm_keys = {k.lower() for k in keys}
        for r in rows:
            mk = (r.get("metric_key") or "").lower()
            if mk in norm_keys:
                return r
        # Fallback by label substring match.
        for r in rows:
            ml = (r.get("metric_label") or "").lower()
            if any(k.replace("_", " ") in ml for k in norm_keys):
                return r
        return None

    metric_row = _find("eps", "earnings_per_share")
    metric_label = "EPS"
    if not metric_row:
        metric_row = _find("net_income", "netincome")
        metric_label = "Net Income"
    if not metric_row:
        metric_row = _find("net_sales", "revenue", "sales")
        metric_label = "Net Sales"
    if not metric_row:
        return [], ""

    by_quarter = metric_row.get("by_quarter") or []
    # Optional revenue row — kept alongside the primary metric so the
    # chart can still annotate revenue regardless of which metric drives
    # the bars.
    rev_row = _find("net_sales", "revenue", "sales") if metric_label != "Net Sales" else None
    rev_bq = rev_row.get("by_quarter") if rev_row else []

    out: list[dict] = []
    for i, q in enumerate(quarters):
        cell = by_quarter[i] if i < len(by_quarter) else {}
        if not isinstance(cell, dict):
            continue
        actual = cell.get("released")
        estimate = cell.get("forecast")
        spread = cell.get("spread_pct")
        # Skip rows with no usable data.
        if not (isinstance(actual, (int, float)) or isinstance(estimate, (int, float))):
            continue
        # Normalise period label to the "Q2 2025" shape the chart parses.
        import re as _re_sh
        m = _re_sh.match(r"\s*(\d{4})\s*Q([1-4])", str(q) or "")
        if m:
            period = f"Q{m.group(2)} {m.group(1)}"
        else:
            m2 = _re_sh.match(r"\s*Q([1-4])\s*(\d{4})", str(q) or "")
            period = f"Q{m2.group(1)} {m2.group(2)}" if m2 else str(q)
        if not isinstance(spread, (int, float)):
            if (isinstance(actual, (int, float)) and isinstance(estimate, (int, float))
                    and estimate not in (0, 0.0)):
                spread = (actual - estimate) / estimate * 100.0
            else:
                spread = None
        row = {
            "period": period,
            "eps_actual": float(actual) if isinstance(actual, (int, float)) else None,
            "eps_estimate": float(estimate) if isinstance(estimate, (int, float)) else None,
            "eps_surprise_pct": float(spread) if isinstance(spread, (int, float)) else None,
        }
        rev_cell = rev_bq[i] if i < len(rev_bq) else {}
        if isinstance(rev_cell, dict):
            row["revenue_actual"] = (rev_cell.get("released")
                                       if isinstance(rev_cell.get("released"), (int, float))
                                       else None)
            row["revenue_estimate"] = (rev_cell.get("forecast")
                                          if isinstance(rev_cell.get("forecast"), (int, float))
                                          else None)
        out.append(row)
    out.sort(key=lambda r: r["period"], reverse=True)
    return out, metric_label


def build_valuation_data(ticker: str, *, analyst_name: str = "Jabal Research",
                            gen_date: str = "",
                            peers_override: Optional[list[dict]] = None,
                            historical_override: Optional[dict] = None,
                            ms_calendar_events: Optional[dict] = None,
                            ) -> ValuationData:
    cv = get_all_fields(ticker)
    # iShares overlay removed (peer table no longer pulls regional ETF).

    profile = cv.get("company_profile")
    pname = (profile.value.get("name") if profile and isinstance(profile.value, dict)
              else ticker)
    # Currency: prefer company_master (the curated listing currency) over
    # the canonical_store profile, because some providers report a
    # company's country-of-record currency (e.g. CNY for Tencent) while
    # the stock actually trades in HKD on HKEX. The deck should reflect
    # the trading currency.
    currency = ""
    try:
        from src.storage.db import load_company as _load_company
        cm = _load_company(ticker) or {}
        currency = (cm.get("currency") or "").strip()
    except Exception:
        pass
    if not currency and profile and isinstance(profile.value, dict):
        currency = profile.value.get("currency") or ""

    # 52w close series. Prefer canonical_store; fall back to the Investing
    # override fetched by the writer (used for yfinance-blocked tickers
    # whose historical_prices field never reaches the canonical store).
    hp = cv.get("historical_prices")
    close_series = []
    if hp and isinstance(hp.value, dict):
        close_series = hp.value.get("close_series") or []
    if not close_series and isinstance(historical_override, dict):
        close_series = historical_override.get("close_series") or []

    # P/E history from valuation_historical. The header reads
    # "P/E MULTIPLE · 5-YEAR RANGE" so we trim the series to the trailing
    # 5 periods. The reference deck shows FY-4 through current; with
    # 8 raw points the chart became unreadable and contradicted the header.
    vh = cv.get("valuation_historical")
    periods, pe_vals = [], []
    if vh and isinstance(vh.value, dict):
        periods = vh.value.get("periods", []) or []
        pe_vals = vh.value.get("pe", []) or []
    if not (periods and pe_vals):
        log.info("[pe-chart] %s: MS /valuation/ empty (periods=%s, pe=%s) — trying Fallback A (Investing)",
                 ticker, len(periods or []), len(pe_vals or []))
    # Fallback A: Investing's pre-computed pe_fy1 / pe_fy2 (derived from
    # last_price / EPS forecasts on the Investing earnings page).
    if not (periods and pe_vals):
        val_fwd = cv.get("valuation_forward")
        fwd = val_fwd.value if val_fwd and isinstance(val_fwd.value, dict) else {}
        pe_fy1 = fwd.get("pe_fy1")
        pe_fy2 = fwd.get("pe_fy2")
        fy1_year = fwd.get("fy1_year")
        fy2_year = fwd.get("fy2_year")
        if isinstance(pe_fy1, (int, float)) and fy1_year:
            periods = [f"FY{fy1_year}"]
            pe_vals = [float(pe_fy1)]
            if isinstance(pe_fy2, (int, float)) and fy2_year:
                periods.append(f"FY{fy2_year}")
                pe_vals.append(float(pe_fy2))

    if not (periods and pe_vals):
        log.info("[pe-chart] %s: Investing valuation_forward empty too — trying Fallback B (synth from MS annual NI)", ticker)
    # Fallback B: synthesize P/E series from MS annual net_income +
    # market_cap. MS publishes 8-year net-income forecasts on /finances/
    # for tickers where the structured /valuation/ table is empty (BKMB,
    # OQEP, most GCC names). We synthesize EPS = NI ÷ shares_outstanding,
    # using shares = market_cap ÷ price. Then P/E = price ÷ EPS for each
    # forecast year. Matches Bloomberg P/E within ~1% on names tested.
    if not (periods and pe_vals):
        # canonical_store stores price + market_cap as SEPARATE scalar
        # fields ("current_price", "market_cap"). The earlier version of
        # this code looked up a non-existent "quote" dict, which made
        # shares_out perpetually None and the synth path never produced
        # bars (the BKMB.OM deck on 2026-05-23 was the visible symptom).
        price_obs = cv.get("current_price")
        mcap_obs  = cv.get("market_cap")
        quote_price = None
        if price_obs is not None:
            pv = price_obs.value
            if isinstance(pv, (int, float)):
                quote_price = float(pv)
            elif isinstance(pv, dict):
                _p = pv.get("price") or pv.get("value")
                quote_price = float(_p) if isinstance(_p, (int, float)) else None
        market_cap = None
        if mcap_obs is not None:
            mv = mcap_obs.value
            if isinstance(mv, (int, float)):
                market_cap = float(mv)
            elif isinstance(mv, dict):
                _m = mv.get("market_cap") or mv.get("value")
                market_cap = float(_m) if isinstance(_m, (int, float)) else None
            # MarketScreener publishes market cap in millions of local
            # currency — same scale-adjustment the snapshot renderer
            # applies. Without it the synth shares-outstanding count is
            # off by a factor of 1e6 and the P/E values blow up to 11
            # MILLIONx, which the sanity-guard then rejects.
            if (mcap_obs.canonical_source or "").lower() == "marketscreener" and market_cap:
                market_cap *= 1_000_000.0
        # Pull MS annual NI series via the raw-observations table —
        # canonical_store carries only metadata, not the full series.
        ms_ann_payload = None
        try:
            from src.storage.db import get_conn as _gc
            import json as _json
            conn = _gc()
            row = conn.execute(
                "SELECT payload_json FROM raw_observations "
                "WHERE provider='marketscreener' AND ticker=? AND field='income_statement_annual' "
                "ORDER BY rowid DESC LIMIT 1",
                (ticker,),
            ).fetchone()
            conn.close()
            if row:
                ms_ann_payload = _json.loads(row["payload_json"])
        except Exception:
            ms_ann_payload = None
        ms_annual = (ms_ann_payload or {}).get("annual") or {}
        ni_series = ms_annual.get("net_income") or []
        ann_periods_raw = ms_annual.get("periods") or []
        # Need market_cap AND price to derive shares; both must be > 0.
        shares_out = None
        if (isinstance(quote_price, (int, float)) and quote_price > 0
            and isinstance(market_cap, (int, float)) and market_cap > 0):
            shares_out = float(market_cap) / float(quote_price)
        log.info("[pe-chart] %s: synth inputs — price=%s market_cap=%s shares_out=%s "
                 "ni_series_len=%d periods=%s ms_ann_payload_present=%s",
                 ticker, quote_price, market_cap, shares_out,
                 len(ni_series or []), ann_periods_raw, bool(ms_ann_payload))
        if shares_out and ni_series and ann_periods_raw:
            synth_periods: list[str] = []
            synth_pe: list[float] = []
            for p, ni in zip(ann_periods_raw, ni_series):
                if isinstance(ni, (int, float)) and ni > 0:
                    eps_synth = float(ni) / shares_out
                    if eps_synth > 0:
                        pe_synth = float(quote_price) / eps_synth
                        if 1 < pe_synth < 200:   # sanity guard
                            synth_periods.append(f"FY{str(p)[-4:]}")
                            synth_pe.append(pe_synth)
            if synth_periods:
                periods = synth_periods
                pe_vals = synth_pe
                log.info("[pe-chart] %s: synth produced %d bars: %s",
                         ticker, len(periods), list(zip(periods, [round(v,1) for v in pe_vals])))
            else:
                log.warning("[pe-chart] %s: synth produced 0 bars (no valid NI rows or sanity-guard rejected all)", ticker)
    if len(periods) > 5 and len(pe_vals) == len(periods):
        periods = periods[-5:]
        pe_vals = pe_vals[-5:]
    current_pe = None
    if pe_vals:
        for v in reversed(pe_vals):
            if isinstance(v, (int, float)):
                current_pe = v
                break

    # Peer comparables — defer to upstream peer-set (peers_override). When
    # nothing is curated, leave the table empty rather than fill it with
    # a regional-ETF proxy (the iShares row was marginal value and added
    # confusion about "what is this peer").
    peers = peers_override or []

    # Rating split + target price + broker actions
    rating_obs = cv.get("rating_split")
    rs = rating_obs.value if rating_obs and isinstance(rating_obs.value, dict) else {}
    rs_normalised = {
        "buy":  int(rs.get("buy",  0) or 0),
        "hold": int(rs.get("hold", 0) or 0),
        "sell": int(rs.get("sell", 0) or 0),
    }
    n_an = int(rs.get("total", sum(rs_normalised.values())) or 0)

    target_obs = cv.get("target_price")
    target_mean = target_high = target_low = None
    if target_obs and isinstance(target_obs.value, dict):
        target_mean = target_obs.value.get("mean")
        target_high = target_obs.value.get("high")
        target_low  = target_obs.value.get("low")
        # Fall back to target_price.n_analysts if rating_split lacked it
        if not n_an:
            n_an = int(target_obs.value.get("n_analysts", 0) or 0)
    # If rating_split has no buy/hold/sell breakdown but we DO know the
    # consensus + total, approximate from the consensus label so the bar
    # isn't 0%/0%/0% on every Saudi/MENA ticker.
    if n_an > 0 and sum(rs_normalised.values()) == 0:
        consensus = (rs.get("consensus") or "").upper() if rs else ""
        if "OUTPERFORM" in consensus or "BUY" in consensus or "ACCUMULATE" in consensus:
            rs_normalised = {"buy": int(n_an * 0.7), "hold": int(n_an * 0.25),
                              "sell": n_an - int(n_an * 0.7) - int(n_an * 0.25)}
        elif "UNDERPERFORM" in consensus or "SELL" in consensus or "REDUCE" in consensus:
            rs_normalised = {"buy": int(n_an * 0.1), "hold": int(n_an * 0.4),
                              "sell": n_an - int(n_an * 0.1) - int(n_an * 0.4)}
        else:  # HOLD / NEUTRAL / unknown
            rs_normalised = {"buy": int(n_an * 0.3), "hold": int(n_an * 0.5),
                              "sell": n_an - int(n_an * 0.3) - int(n_an * 0.5)}

    # Implied % vs last close
    current_price = (cv.get("current_price").value
                       if cv.get("current_price") else None)
    implied = None
    try:
        if target_mean is not None and current_price is not None and float(current_price) > 0:
            implied = (float(target_mean) / float(current_price) - 1.0) * 100
    except (TypeError, ValueError):
        pass

    # Broker actions — canonical_store key 'broker_actions' carries the
    # MS analyst-recommendations list. Each item: {date, headline, source}.
    # The slide renders up to the 3 most recent.
    ba_obs = cv.get("broker_actions")
    broker_actions: list[dict] = []
    if ba_obs and isinstance(ba_obs.value, dict):
        from src.services.render_jabal_snapshot import _normalize_date as _nd
        for item in (ba_obs.value.get("items") or [])[:3]:
            broker_actions.append({
                "date": _nd(item.get("date", "")),
                "text": item.get("headline", ""),
            })

    # Dynamic peer-table label by sector + region. Falls back to
    # "Selected Global Peers" when sector/region can't be inferred.
    industry = (profile.value.get("industry") if profile and isinstance(profile.value, dict) else "") or ""
    sector   = (profile.value.get("sector")   if profile and isinstance(profile.value, dict) else "") or ""
    peer_table_label = _peer_label(ticker, sector, industry)

    # is_bank from curated company_master; controls whether the peer table
    # third multiples column renders EV/EBITDA or P/TBV.
    is_bank_flag = False
    try:
        from src.storage.db import load_company as _load_company_for_bank
        _cm = _load_company_for_bank(ticker) or {}
        is_bank_flag = bool(_cm.get("is_bank"))
    except Exception:
        is_bank_flag = False

    # Earnings surprise history from Investing (already wired into
    # canonical_store as `income_statement_quarterly.surprise_history`).
    # Feeds the new dual-panel earnings-history chart at the bottom of
    # slide 3.
    investing_obs = get_observations_by_provider(ticker, "investing")
    isq = investing_obs.get("income_statement_quarterly") if investing_obs else None
    surprise_history = (isq or {}).get("surprise_history") if isinstance(isq, dict) else None

    # `surprise_history` from Investing can be technically non-empty
    # (list with 10 rows) while carrying no usable data — every row's
    # eps_actual and eps_estimate are null. This is the Investing
    # signature for thinly-covered names like BKMB.OM: the announce-
    # dates are tracked but the values are not. Detect the "list of
    # nulls" case and treat it as empty so the MS fallback below fires.
    def _has_real_surprise_data(rows: list | None) -> bool:
        if not rows or not isinstance(rows, list):
            return False
        for r in rows:
            if not isinstance(r, dict):
                continue
            if (isinstance(r.get("eps_actual"), (int, float))
                    or isinstance(r.get("eps_estimate"), (int, float))
                    or isinstance(r.get("revenue_actual"), (int, float))
                    or isinstance(r.get("revenue_estimate"), (int, float))):
                return True
        return False

    if not _has_real_surprise_data(surprise_history):
        if surprise_history:
            log.info("[earnings-history] %s: Investing returned %d rows but "
                     "all null — falling through to MS fallback",
                     ticker, len(surprise_history))
        surprise_history = None

    # Fallback: when Investing's surprise_history is empty (common on
    # GCC ex-Saudi names — Investing tracks the announce date but has
    # neither actuals nor estimates for BKMB / OQEP / ADNOCDRILL), derive
    # it from MS's /calendar/ quarterly_results table which DOES carry
    # released + forecast per quarter.
    #
    # Two-tier fallback:
    #   1. `ms_calendar_events` parameter (passed by the pipeline). This
    #      is the normal path when fetch_marketscreener_pages succeeded.
    #   2. Direct snapshot read. The pipeline plumbing can drop the
    #      calendar payload silently (status=PARTIAL with HTTP 403 noise
    #      from live attempts). In that case the data DOES exist as a
    #      committed snapshot under data/marketscreener/ — load it
    #      directly so we don't lose the chart over a plumbing miss.
    surprise_metric_label = "EPS"
    if not surprise_history:
        # Tier 1: pipeline-supplied
        derived, derived_label = _surprise_history_from_ms_calendar(ms_calendar_events)
        if derived:
            log.info("[earnings-history] %s: using MS calendar fallback "
                     "(%d quarters of %s data)", ticker, len(derived), derived_label)
            surprise_history = derived
            surprise_metric_label = derived_label or "EPS"
        else:
            # Tier 2: direct snapshot load. Bypasses the pipeline entirely.
            log.info("[earnings-history] %s: ms_calendar_events arrived empty — "
                     "trying direct snapshot read", ticker)
            try:
                snapshot_qr = _load_ms_calendar_from_snapshot(ticker)
            except Exception as exc:
                log.warning("[earnings-history] %s: snapshot read failed: %s", ticker, exc)
                snapshot_qr = None
            if snapshot_qr:
                derived, derived_label = _surprise_history_from_ms_calendar(
                    {"quarterly_results": snapshot_qr}
                )
                if derived:
                    log.info("[earnings-history] %s: snapshot fallback produced "
                             "%d quarters of %s data", ticker, len(derived), derived_label)
                    surprise_history = derived
                    surprise_metric_label = derived_label or "EPS"

    # Optional 5y forward-P/E history if upstream surfaced it; today
    # `valuation_historical` carries the FY-grouped period/pe pair only.
    # The chart falls back to forecast-bars when this is empty.
    pe_history = None
    if vh and isinstance(vh.value, dict):
        pe_history = vh.value.get("forward_pe_history") or None

    return ValuationData(
        company_name=pname,
        close_series=close_series,
        currency=currency,
        pe_periods=periods,
        pe_values=pe_vals,
        pe_current=current_pe,
        peers=peers,
        peer_table_label=peer_table_label,
        rating_split=rs_normalised,
        n_analysts=n_an,
        target_mean=target_mean,
        target_range=(target_low, target_high) if (target_low or target_high) else None,
        target_implied_pct=implied,
        broker_actions=broker_actions,
        sources_line=_sources_line(cv),
        analyst_name=analyst_name,
        gen_date=gen_date or datetime.utcnow().strftime("%d %b %Y"),
        is_bank=is_bank_flag,
        surprise_history=surprise_history,
        surprise_metric_label=surprise_metric_label,
        ticker=ticker,
        pe_history=pe_history,
    )
