"""Earnings preview: PPTX output + sector IV helpers (imported by qa_engine).

Two render paths coexist:

1. Legacy renderer (`_write_preview_pptx_portrait`) — the original
   12-step deck built directly from `ReportPayload`.
2. Jabal renderer (`_write_jabal_preview`) — Stage 2 3-slide deck that
   reads from `canonical_store`. Selected when `JABAL_RENDERER=1` env
   var is set (or `cfg.report.renderer == 'jabal'`).

Legacy stays for backwards-compat; new work should pick Jabal.
"""
from __future__ import annotations
import os
from datetime import datetime
from pathlib import Path

from src.config import cfg, report_output_dir
from src.models.report_payload import ReportPayload
from src.models.step_result import Status, StepResult, StepTimer

STEP = "generate_report"


def _use_jabal_renderer() -> bool:
    """Switch which renderer drives the final PPTX output.

    Order: env override > config flag > default ON.

    The Jabal 3-slide deck is the only user-visible product. Legacy stays
    importable for back-compat but is off by default; set
    `JABAL_RENDERER=0` (or `[report] renderer = "legacy"`) to force it on.
    """
    env_val = os.environ.get("JABAL_RENDERER")
    if env_val is not None:
        return env_val.strip().lower() in {"1", "true", "yes", "on"}
    try:
        return (cfg().get("report", {}) or {}).get("renderer", "jabal") == "jabal"
    except Exception:
        return True


def _write_jabal_preview(payload: ReportPayload, out_path: Path,
                          memo_data: dict | None) -> None:
    """Render the 3-slide Jabal deck via the new render_jabal_* modules,
    reading data from canonical_store.

    Bootstrap step: if canonical_store has no rows for this ticker (e.g.
    fresh DB or first run on a new name), we kick off a daily-cadence
    refresh against the standard provider set so the slides aren't empty.
    """
    from pptx import Presentation
    from pptx.util import Inches as _In
    from src.services.canonical_store import get_all_fields
    from src.services.jabal_design_tokens import PAGE_W_IN, PAGE_H_IN
    from src.services.render_jabal_snapshot import (
        render_snapshot_slide, build_snapshot_data,
    )
    from src.services.render_jabal_thesis import (
        render_thesis_slide, build_thesis_data,
    )
    from src.services.render_jabal_valuation import (
        render_valuation_slide, build_valuation_data,
    )

    ticker = payload.company.ticker
    if not get_all_fields(ticker):
        # Run a quick refresh across the default-on providers so the slides
        # populate. We intentionally exclude Investing.com (Playwright) and
        # IR-PDF (needs curated URL) from the auto-bootstrap.
        import subprocess, sys
        for cadence in ("daily", "weekly", "quarterly"):
            try:
                subprocess.run(
                    [sys.executable, "-m", "scripts.daily_refresh",
                     f"--cadence={cadence}",
                     f"--tickers={ticker}",
                     "--only=yahoo,marketscreener,investing,macro,ishares,commodities"],
                    timeout=180, check=False,
                )
            except Exception:
                continue

    # Stage 2 period defaulting — caller can override via memo_data. When
    # memo_data doesn't carry an explicit period_label / report_date, derive
    # them from payload.memo_computed (preview_quarter_short, next_quarter_label)
    # and payload.yahoo_earnings_date so the snapshot never reads "Earnings
    # Preview · TBA" for a ticker we already have calendar data on.
    mc = getattr(payload, "memo_computed", {}) or {}
    period_label = (memo_data or {}).get("period_label")
    if not period_label:
        # Prefer the explicit "Q<n> <YYYY>" form to the short "<n>Q<YY>"
        nql = mc.get("next_quarter_label") or ""  # e.g. "2026 Q2"
        import re as _rqp
        m = _rqp.search(r"(\d{4})\s*Q(\d)|Q(\d)\s*(\d{4})", nql, _rqp.I)
        if m:
            yr = m.group(1) or m.group(4)
            qn = m.group(2) or m.group(3)
            period_label = f"Q{qn} {yr} Earnings Preview"
        else:
            period_label = mc.get("preview_quarter_label") or "Earnings Preview"

    report_date = (memo_data or {}).get("report_date")
    if not report_date:
        # MS /calendar/ block exposes next_expected_earnings_date (ISO date).
        ms_cal = getattr(payload, "ms_calendar_events", None) or {}
        for path in (
            ms_cal.get("next_expected_earnings_date") if isinstance(ms_cal, dict) else None,
            ms_cal.get("next_expected_earnings_label") if isinstance(ms_cal, dict) else None,
            getattr(payload, "yahoo_earnings_date", None),
            mc.get("next_earnings_date"),
            mc.get("yahoo_earnings_date"),
        ):
            if path:
                report_date = str(path)
                break
        report_date = report_date or "TBA"

    # Load curated peer tickers (from company_master.peer_group) and
    # enrich them with yfinance — drives slide 3's peer table.
    peer_rows: list[dict] = []
    try:
        from src.storage.db import load_company
        from src.services.fetch_peers import fetch_peer_rows
        company_row = load_company(ticker) or {}
        peer_tickers = company_row.get("peer_group") or []
        if peer_tickers:
            peer_rows = fetch_peer_rows(peer_tickers)
    except Exception as exc:
        import logging
        logging.getLogger(__name__).warning("peer enrichment failed for %s: %s", ticker, exc)

    snap      = build_snapshot_data(ticker, period_label=period_label,
                                       report_date=report_date,
                                       ms_price_performance=getattr(payload, "ms_price_performance", None))
    is_bank = bool(getattr(payload.company, "is_bank", False))
    thesis    = build_thesis_data(ticker,
                                     quarterly=getattr(payload, "quarterly_actuals", None) or [],
                                     is_bank=is_bank,
                                     ms_quarterly_forecasts=getattr(payload, "ms_quarterly_forecasts", None))
    valuation = build_valuation_data(ticker, peers_override=peer_rows or None)

    prs = Presentation()
    prs.slide_width  = _In(PAGE_W_IN)
    prs.slide_height = _In(PAGE_H_IN)
    render_snapshot_slide(prs, snap)
    render_thesis_slide(prs, thesis)
    render_valuation_slide(prs, valuation)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

# Minimum character length for IV paragraphs before using fallback (with Recent Context we accept shorter LLM output)
MIN_IV_LEN_WITH_RECENT_CONTEXT = 20
MIN_IV_LEN_DEFAULT = 40
IV_STYLE_DEFAULT = "balanced"
IV_STYLE_ALLOWED = {"balanced", "tactical", "conservative"}



def _fmt_num(val, in_millions: bool = False) -> str:
    if val is None:
        return "—"
    if isinstance(val, (int, float)):
        if in_millions and val >= 1e3:
            return f"{val / 1e3:,.2f}bn"
        if val >= 1e9:
            return f"{val / 1e9:.2f}B"
        if val >= 1e6:
            return f"{val / 1e6:,.0f}M"
        if abs(val) < 1e-6 and val != 0:
            return f"{val:.4f}"
        return f"{val:,.2f}" if val != int(val) else f"{int(val):,}"
    return str(val)


def _fmt_pct(val, signed: bool = False) -> str:
    if val is None:
        return "—"
    if isinstance(val, (int, float)):
        return f"{val:+.1f}%" if signed else f"{val}%"
    return str(val)


def _field_display(f, default: str = "—"):
    """Return display value only when field status allows rendering; else default."""
    if not f or not isinstance(f, dict):
        return default
    if f.get("status") not in ("pass", "stale", "estimated", "manually_entered"):
        return default
    v = f.get("display_value") if f.get("display_value") is not None else f.get("value")
    return v if v is not None else default


def _company_attr(company, key: str, default: str = ""):
    """Read sector/industry/etc. from company whether it's a model instance or a dict (e.g. after serialization)."""
    if company is None:
        return default
    if isinstance(company, dict):
        return (company.get(key) or default) if isinstance(default, str) else company.get(key, default)
    return (getattr(company, key, None) or default) if isinstance(default, str) else getattr(company, key, default)


def _sector_operating_kpis_and_what_matters(company) -> tuple[list[str], list[str], str]:
    """
    Return (operating_metrics_kpis[4], what_matters_bullets[5], fallback_para2_snippet).
    fallback_para2 is publishable analyst prose only (no "Focus on...", "Do not use..." instructions).
    """
    sector = (_company_attr(company, "sector", "") or "").strip().lower()
    industry = (_company_attr(company, "industry", "") or "").strip().lower()
    ind = industry or sector
    is_bank = bool(_company_attr(company, "is_bank", False))

    if is_bank:
        kpis = ["Loans", "Deposits", "NIM", "Cost of Risk"]
        matters = ["Loan / financing growth", "NIM / margin", "Asset quality", "Funding mix", "Capital return"]
        p2 = "For banks, the story usually turns on NIM, loan growth, and asset quality. Earnings quality—recurring versus one-offs—and any weakness versus consensus or deterioration in asset quality are key for the multiple."
        return kpis, matters, p2

    if "oil" in ind or "gas" in ind or "energy" in sector or "exploration" in ind or "petroleum" in ind:
        kpis = ["Production volumes", "Realized oil/gas prices", "Lifting costs", "Capex / project ramp-up"]
        matters = ["Production volumes", "Realized oil/gas prices", "Lifting costs", "Reserve replacement / field startup", "Capex and project ramp-up"]
        p2 = "For oil and gas, the narrative typically turns on production volumes, realized prices, and lifting costs; reserve replacement and field startup impact also matter, and capex and project ramp-up often drive the story."
        return kpis, matters, p2

    if "telecom" in ind or "communication" in sector or "communication" in ind:
        kpis = ["Subscribers", "ARPU", "Churn", "Capex intensity"]
        matters = ["Subscriber additions", "ARPU trend", "Churn", "Capex intensity", "India wireless competition" if "india" in ((_company_attr(company, "country", "") or "").lower()) else "Wireless competition", "Enterprise / data centre contribution"]
        p2 = "For telecoms and communication equipment, subscriber trends, ARPU, churn, and capex intensity are central where relevant; enterprise and product-cycle dynamics often drive the story."
        return kpis, matters[:5], p2

    if "technology" in sector or "software" in ind or "semiconductor" in ind or "equipment" in ind:
        kpis = ["Revenue growth", "Margin", "Guidance", "Key product metrics"]
        matters = ["Revenue mix and growth", "Margin and profitability", "Guidance", "Product cycles", "Competitive dynamics"]
        p2 = "For technology and communication equipment names, the narrative typically turns on revenue mix, margins, and guidance; product cycles and competitive dynamics often drive the stock."
        return kpis, matters[:5], p2

    if "industrial" in sector or "capital good" in ind or "aerospace" in ind or "machinery" in ind:
        kpis = ["Orders / backlog", "Utilization", "Pricing", "Guidance"]
        matters = ["Demand and orders", "Backlog / utilization", "Margin and pricing", "Guidance", "Key metrics"]
        p2 = "For industrials, demand, orders, backlog, and utilization drive the story; margin and pricing matter, and guidance and key metrics often move the stock."
        return kpis, matters, p2

    if "internet" in ind or "e-commerce" in ind or "retail" in ind:
        kpis = ["GMV", "Cloud revenue growth", "International commerce", "Customer management revenue"]
        matters = ["GMV and engagement", "Margin and pricing", "Guidance", "Key metrics"]
        p2 = "Sector operating metrics and headline results versus consensus are key; guidance and main metrics typically drive the stock."
        return kpis, matters[:5], p2

    if "mining" in ind or "metals" in ind:
        kpis = ["Production / throughput", "Commodity prices", "Costs", "Guidance"]
        matters = ["Production and sales volume", "Commodity price realizations", "Cost and margin", "Guidance", "Key metrics"]
        p2 = "For metals and mining, the narrative typically turns on production, realized commodity prices, and costs; guidance and key operating metrics often drive the stock."
        return kpis, matters[:5], p2

    if "chem" in ind or "material" in ind or "material" in sector:
        kpis = ["Volume", "Realized price", "Utilization", "Feedstock spread"]
        matters = ["Volume and realized price", "Utilization", "Feedstock spread", "Guidance", "Key metrics"]
        p2 = "Volume, realized price, utilization, and feedstock spread are the main levers; guidance and key metrics drive the story."
        return kpis, matters[:5], p2

    if "real estate" in sector or "real estate" in ind or "reit" in ind or "property" in ind:
        kpis = ["Occupancy", "Rental rates", "Recurring income mix", "Development pipeline"]
        matters = ["Occupancy trend", "Rental and sales prices", "Recurring vs development income", "Development pipeline and pre-sales", "Leverage and cost of debt"]
        p2 = "For real estate, the narrative typically turns on occupancy, rental rates, and the mix of recurring versus development income; pipeline execution, pre-sales, and leverage drive the stock."
        return kpis, matters[:5], p2

    if "insurance" in sector or "insurance" in ind:
        kpis = ["Gross written premium", "Combined ratio", "Investment yield", "Solvency"]
        matters = ["Premium growth", "Underwriting margin and combined ratio", "Investment income and asset mix", "Claims experience", "Capital and solvency"]
        p2 = "For insurers, the story centers on premium growth, the combined ratio, and investment yield; claims experience and solvency drive the multiple."
        return kpis, matters[:5], p2

    if "financial" in sector or "capital markets" in ind or "asset management" in ind or "investment" in ind:
        kpis = ["AUM / assets", "Fee rate", "Cost/income", "Capital return"]
        matters = ["Asset growth", "Fee rate and revenue mix", "Cost efficiency", "Asset quality / risk-weighted assets", "Capital return"]
        p2 = "For diversified financials, the narrative turns on asset growth, fee rate, and cost/income; risk posture and capital return drive the valuation."
        return kpis, matters[:5], p2

    if "utilities" in sector or "utility" in ind or "water" in ind or "electric" in ind or ("gas" in ind and "oil" not in ind):
        kpis = ["Regulated asset base", "Tariff / allowed return", "Volumes", "Capex"]
        matters = ["Tariff and allowed return", "Regulated asset base growth", "Demand / volumes", "Capex delivery", "Guidance"]
        p2 = "For utilities, the story typically turns on regulated asset base growth, tariffs, and volumes; capex execution and any guidance changes drive the stock."
        return kpis, matters[:5], p2

    if "consumer staples" in sector or "food" in ind or "beverage" in ind or "household" in ind or "personal product" in ind:
        kpis = ["Volume", "Price/mix", "Gross margin", "Guidance"]
        matters = ["Volume and price/mix", "Gross margin and input costs", "Market share", "Marketing / promotional intensity", "Guidance"]
        p2 = "For consumer staples, the narrative is driven by volume and price/mix; input cost and promotional intensity set margins, and guidance anchors the multiple."
        return kpis, matters[:5], p2

    if "consumer discretionary" in sector or "apparel" in ind or "auto" in ind or "restaurant" in ind or "lodging" in ind or "specialty retail" in ind:
        kpis = ["Same-store sales / volumes", "Price/mix", "Margin", "Guidance"]
        matters = ["Same-store sales and unit growth", "Price/mix and discounting", "Gross and operating margin", "Inventory position", "Guidance"]
        p2 = "For consumer discretionary, same-store sales, price/mix, and margin are the core drivers; inventory discipline and guidance shape the stock."
        return kpis, matters[:5], p2

    if "healthcare" in sector or "pharma" in ind or "biotech" in ind or "medical" in ind or "drug" in ind:
        kpis = ["Revenue by franchise", "Gross margin", "R&D progress", "Guidance"]
        matters = ["Franchise performance", "Pricing and volume", "Gross and operating margin", "Pipeline / R&D milestones", "Guidance"]
        p2 = "For healthcare names, franchise revenue trends, margins, and pipeline milestones dominate; pricing pressure and guidance anchor sentiment."
        return kpis, matters[:5], p2

    # Default: labeled rows for manual entry. Kept generic on purpose so the report
    # makes clear the section is a placeholder rather than claiming sector-specific focus.
    kpis = ["Key metric 1", "Key metric 2", "Key metric 3", "Key metric 4"]
    matters = ["Headline vs consensus", "Margin and pricing", "Guidance", "Key metrics"]
    p2 = "This quarter, sector operating metrics and headline results versus consensus matter most; earnings quality—whether a beat or miss is recurring or one-off—and guidance drive the narrative."
    return kpis, matters[:5], p2


def _iv_fallback_style() -> str:
    """
    Fallback IV style selector.
    Priority: env IV_FALLBACK_STYLE -> config report.iv_fallback_style -> default balanced.
    """
    style = (os.environ.get("IV_FALLBACK_STYLE") or "").strip().lower()
    if not style:
        try:
            style = str(cfg().get("report", {}).get("iv_fallback_style", "")).strip().lower()
        except Exception:
            style = ""
    if style not in IV_STYLE_ALLOWED:
        style = IV_STYLE_DEFAULT
    return style


def _build_analytical_iv_paragraph_1(
    company_name: str,
    preview_short: str,
    rec: str,
    an_str: str,
    price: float | None,
    spread: float | None,
    rev_surprise: float | None,
    eps_surprise: float | None,
    memo: dict,
    _fmt_pct,
    _fmt_num,
    style: str = IV_STYLE_DEFAULT,
) -> str:
    """
    Build a single analytical paragraph for the Investment View fallback.
    Interprets consensus, surprise history, and key preview rather than listing data points.
    """
    sentences = []

    # Opening: frame the setup
    if style == "tactical":
        sentences.append(f"Into the {preview_short} print, {company_name} screens as a tactical setup.")
    elif style == "conservative":
        sentences.append(f"Into {preview_short}, {company_name} has a constructive but not low-risk setup.")
    else:
        sentences.append(f"{company_name} reports {preview_short}.")

    # Street view and what underpins it
    if rec and rec != "—":
        line = f"The street has a {rec} rating"
        if an_str:
            line += f" ({an_str})"
        if price is not None:
            line += f", with the average target at {_fmt_num(price)}"
            if spread is not None:
                line += f", implying {_fmt_pct(spread, signed=True)} upside"
        line += "."
        sentences.append(line)
        # Interpret surprise history
        if rev_surprise is not None and eps_surprise is not None:
            rev_beat = rev_surprise > 0
            eps_beat = eps_surprise > 0
            if rev_beat and not eps_beat:
                if style == "tactical":
                    sentences.append(
                        f"Revenue has tended to beat (avg {_fmt_pct(rev_surprise, signed=True)}) while EPS has lagged ({_fmt_pct(eps_surprise, signed=True)}); "
                        "the immediate trigger is whether top-line resilience converts into cleaner earnings."
                    )
                elif style == "conservative":
                    sentences.append(
                        f"Revenue has tended to beat consensus (avg {_fmt_pct(rev_surprise, signed=True)}), "
                        f"while EPS has lagged ({_fmt_pct(eps_surprise, signed=True)}), so earnings conversion remains the main risk into the quarter."
                    )
                else:
                    sentences.append(
                        f"Revenue has tended to beat consensus (avg {_fmt_pct(rev_surprise, signed=True)}), "
                        f"while EPS has lagged ({_fmt_pct(eps_surprise, signed=True)}); the story into the print hinges on whether top-line strength can translate into earnings delivery."
                    )
            elif not rev_beat and eps_beat:
                sentences.append(
                    f"EPS has run ahead of consensus (avg {_fmt_pct(eps_surprise, signed=True)}), though revenue has been softer ({_fmt_pct(rev_surprise, signed=True)}); the focus will be on sustainability of margins and guidance."
                )
            elif rev_beat and eps_beat:
                sentences.append(
                    f"Both revenue and EPS have tended to beat (revenue avg {_fmt_pct(rev_surprise, signed=True)}, EPS {_fmt_pct(eps_surprise, signed=True)}), which supports the constructive setup but raises the bar for this quarter."
                )
            else:
                sentences.append(
                    f"Versus consensus, revenue has averaged {_fmt_pct(rev_surprise, signed=True)} and EPS {_fmt_pct(eps_surprise, signed=True)}; the quarter will need to show improvement or a clear path to it for the rating to hold."
                )
        elif rev_surprise is not None:
            sentences.append(
                f"Revenue versus consensus has averaged {_fmt_pct(rev_surprise, signed=True)}; "
                + ("that consistency supports the constructive view." if rev_surprise > 0 else "delivery this quarter will be important for confidence.")
            )
        elif eps_surprise is not None:
            sentences.append(
                f"EPS surprise has averaged {_fmt_pct(eps_surprise, signed=True)}; "
                + ("earnings delivery has underpinned the rating." if eps_surprise > 0 else "the market will be looking for better earnings consistency.")
            )

    # Key preview: tougher comp / context for the quarter
    calendar_prior = memo.get("calendar_prior_quarter_released") or {}
    calendar_same_ly = memo.get("calendar_same_q_prior_yr_released") or {}
    has_prior = (calendar_prior.get("net_sales") is not None) or (memo.get("prior_quarter_actual_revenue") is not None)
    has_same_ly = (calendar_same_ly.get("net_sales") is not None) or (memo.get("same_quarter_prior_year_revenue") is not None)
    qoq_rev = memo.get("qoq_revenue_pct")
    yoy_rev = memo.get("yoy_revenue_pct_table")
    if (qoq_rev is not None and has_prior) or (yoy_rev is not None and has_same_ly):
        q_part = _fmt_pct(qoq_rev, signed=True) if (qoq_rev is not None and has_prior) else None
        y_part = _fmt_pct(yoy_rev, signed=True) if (yoy_rev is not None and has_same_ly) else None
        if q_part is not None or y_part is not None:
            bits = []
            if q_part is not None:
                bits.append(f"QoQ {q_part}")
            if y_part is not None:
                bits.append(f"YoY {y_part}")
            preview_phrase = " and ".join(bits)
            if style == "tactical":
                sentences.append(
                    f"Key preview points to {preview_phrase}; with tougher comps, an in-line outcome may be enough only if margin/cost delivery is clean."
                )
            elif style == "conservative":
                sentences.append(
                    f"Key preview points to {preview_phrase}—a tougher comparison that raises execution risk; "
                    "the quarter needs a credible delivery path to preserve confidence."
                )
            else:
                sentences.append(
                    f"The key preview points to {preview_phrase}—a tougher comparison; "
                    "the focus will be on whether the company can meet or beat the bar and sustain the narrative."
                )

    # Expectations into the print
    if spread is not None:
        if spread > 0:
            if style == "tactical":
                sentences.append("Expectations look supportive; an in-line or better print likely keeps near-term positioning constructive.")
            elif style == "conservative":
                sentences.append("Expectations look supportive, but the reaction still depends on earnings quality and guidance credibility.")
            else:
                sentences.append("Expectations into the print look supportive; an in-line or better outcome would likely be well received.")
        elif spread < 0:
            sentences.append("Expectations look demanding; the stock may need a clear beat or raise to re-rate.")
        else:
            sentences.append("Expectations into the print look balanced.")
    else:
        sentences.append("Expectations into the print look balanced.")

    return " ".join(sentences)



def _iv_text_and_watch(payload: ReportPayload, memo_data: dict | None, iv_style: str) -> tuple[str, list[str]]:
    """LLM IV if long enough; else analytical + sector p2 + optional recent-coverage snippet."""
    c, memo = payload.company, payload.memo_computed or {}
    sections = (memo_data or {}).get("pptx_sections") or {}
    if isinstance(sections, dict):
        thesis = (sections.get("investment_thesis") or "").strip()
        wtw = sections.get("what_to_watch") if isinstance(sections.get("what_to_watch"), list) else []
        wtw = [str(x).strip() for x in (wtw or []) if str(x).strip()]
        if thesis:
            return thesis, (wtw[:4] if wtw else _sector_operating_kpis_and_what_matters(c)[1])
    ns = getattr(payload, "news_summary", None)
    min_len = MIN_IV_LEN_WITH_RECENT_CONTEXT if (ns and getattr(ns, "referenced_articles", None)) else MIN_IV_LEN_DEFAULT
    p1 = (getattr(ns, "investment_view_paragraph_1", "") or "").strip() if ns else ""
    p2 = (getattr(ns, "investment_view_paragraph_2", "") or "").strip() if ns else ""
    _, matters, p2_fb = _sector_operating_kpis_and_what_matters(c)
    if len(p1) >= min_len and len(p2) >= min_len:
        return f"{p1} {p2}".strip(), matters
    header = (memo_data or {}).get("header") or {}
    rec = _field_display(header.get("recommendation"), "—")
    n_an = _field_display(header.get("analyst_count"))
    try:
        n_an = int(n_an) if n_an not in (None, "—", "") else None
    except (TypeError, ValueError):
        n_an = None
    an_str = f"{n_an} analysts" if isinstance(n_an, int) and n_an else ""
    price = _field_display(header.get("average_target_price"))
    try:
        price = float(price) if price not in (None, "—", "") else None
    except (TypeError, ValueError):
        price = None
    spread = _field_display(header.get("upside_pct"))
    try:
        spread = float(spread) if spread not in (None, "—", "") else None
    except (TypeError, ValueError):
        spread = None
    preview_short = (memo_data or {}).get("preview_short") or memo.get("preview_quarter_short") or f"{(datetime.now().month - 1) // 3 + 1}Q{datetime.now().strftime('%y')}"
    company_name = getattr(c, "company_name", None) or _company_attr(c, "company_name", "")
    fb1 = _build_analytical_iv_paragraph_1(
        company_name=company_name,
        preview_short=preview_short,
        rec=rec,
        an_str=an_str,
        price=price,
        spread=spread,
        rev_surprise=memo.get("avg_revenue_surprise_pct"),
        eps_surprise=memo.get("avg_eps_surprise_pct"),
        memo=memo,
        _fmt_pct=_fmt_pct,
        _fmt_num=_fmt_num,
        style=iv_style,
    )
    # Headlines are surfaced in the dedicated "Recent Headlines" sidebar
    # (slide 2). Previously a single news fact was stitched mid-thesis here,
    # which produced run-on sentences like
    #   "...Argaam Volume, realized price, utilization..."
    # because the fact ended without punctuation and concatenated into the
    # "Volume, realized price..." key-focus list. Keeping the thesis prose
    # clean and routing all news to the sidebar fixes the class of bug.
    # Expand fallback IV so exec summary is more comprehensive even if Gemini fails.
    focus = ""
    try:
        focus_bits = [m for m in (matters or []) if m][:4]
        if focus_bits:
            focus = " Key focus areas include " + ", ".join(focus_bits) + "."
    except Exception:
        focus = ""
    return f"{fb1} {p2_fb}{focus}".strip(), matters



def _write_preview_pptx_portrait(
    payload: ReportPayload,
    path: Path,
    memo_data: dict | None,
    iv_text: str,
    watch: list[str],
    quality_flags: list[str] | None = None,
) -> None:
    """Portrait-oriented (7.5 × 13.33 in) PPTX.

    Slide 1 (cover), 2 (executive summary), and 3 (financial snapshot) are
    rendered from a typed `ReportContext` produced by
    `build_report_context.build()`. Slide 4 (disclosures) is static and
    rendered inline below — no per-deck values to thread.

    The legacy implementation computed every slide's values inline in this
    function; that duplicated computation is gone, eliminating the entire
    class of "cover says X, thesis says Y" drift bugs.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # Colour palette used by slide 4 (the only slide still drawn inline).
    DARK = RGBColor(0x0D, 0x11, 0x17)
    GOLD = RGBColor(0xC9, 0xA2, 0x27)
    LIGHT = RGBColor(0xE6, 0xED, 0xF3)
    MUTED = RGBColor(0x8B, 0x94, 0x9E)

    # Pptx primitive closures shared with the per-slide renderers. These are
    # passed in via kwargs so each renderer module is standalone (does not
    # import pptx itself except via these helpers + chart_builders).
    def tx(sl, x, y, w, h, t, *, sz=12, bold=False, rgb=RGBColor(0, 0, 0),
           al=PP_ALIGN.LEFT, word_wrap=True, line_spacing=None,
           hyperlink: str | None = None):
        b = sl.shapes.add_textbox(x, y, w, h)
        tf = b.text_frame
        tf.clear()
        tf.word_wrap = word_wrap
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        text = "" if t is None else str(t)
        lines = text.split("\n") if text else [""]

        def _set(para, ln):
            para.text = ln
            para.alignment = al
            if line_spacing is not None:
                try:
                    para.line_spacing = line_spacing
                except Exception:
                    pass
            try:
                para.space_after = Pt(0)
                para.space_before = Pt(0)
            except Exception:
                pass
            if para.runs:
                run = para.runs[0]
                run.font.name = "Arial"
                run.font.size = Pt(sz)
                run.font.bold = bold
                run.font.color.rgb = rgb
                # When a hyperlink is supplied, attach it to every run so the
                # whole textbox is clickable. PPTX hyperlinks live on runs,
                # not on the text frame, so multi-line links need this for
                # each paragraph's first run.
                if hyperlink:
                    try:
                        run.hyperlink.address = hyperlink
                    except Exception:
                        pass

        _set(tf.paragraphs[0], lines[0])
        for ln in lines[1:]:
            _set(tf.add_paragraph(), ln)

    def rect(sl, x, y, w, h, fill, line=None, lw=1.0):
        sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if line is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = line
            sh.line.width = Pt(lw)
        return sh

    # ── Build typed ReportContext, render slides 1-3 ──────────
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(13.33)
    W = prs.slide_width
    blank = prs.slide_layouts[6]

    from src.services.build_report_context import build as _build_ctx
    from src.services.render_cover import render as _render_cover
    from src.services.render_summary import render as _render_summary
    from src.services.render_snapshot import render as _render_snapshot
    from src.services.render_income_evolution import render as _render_income_evolution
    from src.services.render_ratings import render as _render_ratings
    from src.services.render_peers import render as _render_peers
    from src.services.render_price_action import render as _render_price_action
    _ctx = _build_ctx(
        payload, memo_data,
        iv_text=iv_text, watch=watch,
        quality_flags=list(quality_flags or []),
    )
    _render_cover(prs, blank, _ctx.cover, tx=tx, rect=rect)
    _render_summary(prs, blank, _ctx.summary, tx=tx, rect=rect)
    _render_snapshot(
        prs, blank, _ctx.snapshot, tx=tx, rect=rect,
        quality_flags=_ctx.quality_flags,
    )

    # MS-extras slides — each renders only when its `has_data` flag is set.
    # Keeping them gated here (rather than inside the renderer) makes the
    # slide order explicit and avoids surfacing partially-empty slides on
    # tickers MarketScreener has thin coverage for.
    if _ctx.income_evolution and _ctx.income_evolution.has_data:
        _render_income_evolution(
            prs, blank, _ctx.income_evolution,
            tx=tx, rect=rect,
            company_name=_ctx.company_name,
        )
    if _ctx.ratings and _ctx.ratings.has_data:
        _render_ratings(
            prs, blank, _ctx.ratings,
            tx=tx, rect=rect,
            company_name=_ctx.company_name,
        )
    if _ctx.sector and _ctx.sector.has_data:
        _render_peers(
            prs, blank, _ctx.sector,
            tx=tx, rect=rect,
            company_name=_ctx.company_name,
        )
    if _ctx.price_action and _ctx.price_action.has_data:
        _render_price_action(
            prs, blank, _ctx.price_action,
            tx=tx, rect=rect,
            company_name=_ctx.company_name,
            currency=_ctx.currency,
        )

    # ── Slide 4: Important Disclosures (dark, portrait) ───────
    s4 = prs.slides.add_slide(blank)
    rect(s4, 0, 0, W, prs.slide_height, DARK)
    tx(s4, Inches(0), Inches(1.0), W, Inches(0.6), "Important Disclosures", sz=28, bold=True, rgb=LIGHT, al=PP_ALIGN.CENTER)
    rect(s4, Inches(2.7), Inches(1.6), Inches(2.1), Inches(0.05), GOLD)
    disclosures = (
        "This document is provided for informational purposes only and does not constitute an offer, "
        "solicitation, or recommendation to buy or sell any security. The information contained herein "
        "is based on sources believed to be reliable, but no representation or warranty, express or "
        "implied, is made regarding its accuracy, completeness, or timeliness.\n\n"
        "All financial data, estimates, and projections are derived from publicly available sources "
        "including MarketScreener and Yahoo Finance, supplemented by AI-generated qualitative analysis. "
        "Past performance is not indicative of future results. Investors should conduct their own due "
        "diligence and consult with a qualified financial advisor before making investment decisions.\n\n"
        "This report does not take into account the specific investment objectives, financial situation, "
        "or particular needs of any individual investor. The securities discussed may not be suitable for "
        "all investors. Investing involves risks, including the possible loss of principal."
    )
    tx(s4, Inches(0.8), Inches(2.0), Inches(5.9), Inches(5.0), disclosures, sz=11, rgb=MUTED, line_spacing=1.3)
    gen_ts = datetime.now().strftime("%d %B %Y at %H:%M UTC")

    # Build the Data Sources line dynamically from what actually fed the deck.
    # When Bloomberg overrode MS+Yahoo, surface that prominently so any
    # reviewer can see at a glance which dataset drove the numbers.
    _used = []
    _act = (_ctx.snapshot.table.actuals_source or "").strip()
    _est = (_ctx.snapshot.table.estimates_source or "").strip()
    if _act == "Bloomberg" or _est == "Bloomberg":
        _used.append("Bloomberg")
    if "MarketScreener" in (_act, _est) or _ctx.summary.income_chart and (_ctx.summary.income_chart.source_label or "") == "MarketScreener":
        _used.append("MarketScreener")
    if "Yahoo" in _act or "Yahoo" in _est:
        _used.append("Yahoo Finance")
    # Google Gemini is included whenever we shipped LLM-generated thesis text.
    if (_ctx.summary.thesis_source or "") == "gemini":
        _used.append("Google Gemini")
    # Defensive default: always show at least one source string.
    if not _used:
        _used = ["MarketScreener", "Yahoo Finance"]
    tx(s4, Inches(0.8), Inches(7.5), Inches(5.9), Inches(0.5),
       f"Data Sources: {', '.join(_used)}\nGenerated: {gen_ts}",
       sz=9, rgb=RGBColor(0x60, 0x66, 0x70), al=PP_ALIGN.CENTER)
    tx(s4, Inches(0), Inches(12.8), W, Inches(0.3),
       f"\u00a9 {datetime.now().year} Earnings Research  |  All rights reserved",
       sz=9, rgb=RGBColor(0x60, 0x66, 0x70), al=PP_ALIGN.CENTER)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def run(payload: ReportPayload, memo_data: dict | None = None, qa_audit: dict | None = None, data_warnings: list[str] | None = None) -> StepResult:
    with StepTimer() as t:
        try:
            iv_style = _iv_fallback_style()
            ticker = payload.company.ticker
            out_dir = report_output_dir()
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            suffix = f"{ticker}_{ts}_earnings_preview.pptx"
            out_path = out_dir / suffix
            iv_text, watch = _iv_text_and_watch(payload, memo_data, iv_style)
            quality_flags: list[str] = []
            if qa_audit:
                # The legacy flags are too granular and overlap — for tickers
                # where MS was rate-limited or unreachable mid-run, the deck
                # rendered "MS entity mismatch suppressed; MS suppressed:
                # missing current data; MS suppressed: entity mismatch" which
                # made it look like a data-quality problem with the ticker
                # rather than a transient outage. Collapse those into one
                # honest message based on whether MS data was FETCHED at all
                # (no lineage = unavailable) vs FETCHED-BUT-REJECTED (lineage
                # present but entity check failed).
                no_lineage = qa_audit.get(
                    "ms_section_suppressed_due_to_missing_current_data", False
                )
                entity_mismatch = (
                    qa_audit.get("ms_section_suppressed_due_to_entity_mismatch", False)
                    or not qa_audit.get("payload_entity_match", True)
                )
                contamination = qa_audit.get("ms_section_suppressed_due_to_contamination", False)
                if contamination:
                    quality_flags.append("MarketScreener data flagged as cross-company duplicate")
                elif no_lineage:
                    # Most common cause: MS was unavailable (rate limit / network /
                    # captcha) when this run executed. NOT a data-correctness issue.
                    quality_flags.append("MarketScreener data unavailable for this run — falling back to Yahoo Finance")
                elif entity_mismatch:
                    # MS data WAS fetched but lineage failed entity validation —
                    # genuine wrong-entity case that needs the ticker mapping fixed.
                    quality_flags.append("MarketScreener data rejected (entity mismatch)")
                if qa_audit.get("reused_default_payload_detected"):
                    quality_flags.append("Default payload reused")
            # Add automated data validation warnings
            if data_warnings:
                quality_flags.extend(data_warnings)
            if _use_jabal_renderer():
                # New Stage-2 3-slide deck — reads exclusively from
                # canonical_store. Bootstraps a quick refresh if the
                # store has no rows for this ticker.
                _write_jabal_preview(payload, out_path, memo_data)
                source_tag = "pptx-jabal"
            else:
                _write_preview_pptx_portrait(
                    payload, out_path, memo_data, iv_text, watch,
                    quality_flags or None,
                )
                source_tag = "pptx"
            return StepResult(step_name=STEP, status=Status.SUCCESS, source=source_tag, message=f"Report saved → {out_path}", data=str(out_path), elapsed_seconds=t.elapsed)
        except Exception as exc:
            return StepResult(step_name=STEP, status=Status.FAILED, source="pptx", message="Report generation failed", error_detail=str(exc), elapsed_seconds=t.elapsed)
