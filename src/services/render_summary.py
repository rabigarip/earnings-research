"""
Render slide 2 (Executive Summary) of the earnings preview deck.

Consumes only `SummaryData`; never reaches into `payload` or `memo_data`.
Layout matches the legacy portrait deck pixel-for-pixel so the visual
diff against the gold-standard PDF is zero on the regression suite.

Sections, top to bottom:
- Header strip: company | period
- "Executive Summary" title with gold accent rule
- Investment Thesis box (LLM or analytical fallback)
- Income Statement Evolution + P/E charts (side-by-side row)
- Source chip (top-right of chart row)
- Key Expectations: 3 cards with sign-aware delta chips
- What to Watch: numbered bullets (omitted if empty per editorial policy)
- Catalysts & Risks: green / red boxes (omitted if empty)
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from src.models.report_context import SummaryData
from src.services.chart_builders import build_pe_chart, build_revenue_ni_chart


# Colour palette — kept in sync with `render_cover.py` and the legacy renderer.
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1F, 0x23, 0x28)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
RED = RGBColor(0xCF, 0x22, 0x22)
THESIS_BG = RGBColor(0xFA, 0xF8, 0xF3)
THESIS_BORDER = RGBColor(0xDB, 0xE0, 0xE6)
CATALYST_BG = RGBColor(0xF0, 0xF9, 0xF0)
CATALYST_BAR = RGBColor(0x1A, 0x7F, 0x37)
RISK_BG = RGBColor(0xFE, 0xF0, 0xF0)
RISK_BAR = RGBColor(0xCF, 0x22, 0x22)


def _delta_color(v: float | None) -> RGBColor:
    if v is None:
        return MUTED
    return GREEN if v >= 0 else RED


def render(prs, blank_layout, summary: SummaryData, *, tx, rect) -> None:
    """Render the Executive Summary slide.

    `tx` and `rect` are the textbox / rect helpers from the parent module
    (closures with the right styling defaults). `prs` is portrait-sized.
    """
    W = prs.slide_width

    s2 = prs.slides.add_slide(blank_layout)
    rect(s2, 0, 0, W, prs.slide_height, WHITE)

    # ── Header strip ──
    head_text = (
        f"{summary.company_name} | {summary.period_label}"
        if summary.company_name and summary.period_label
        else (summary.company_name or summary.period_label or "")
    )
    tx(s2, Inches(0.6), Inches(0.4), Inches(6.3), Inches(0.3),
       head_text, sz=12, bold=True, rgb=BLACK)
    tx(s2, Inches(0.6), Inches(0.85), Inches(6), Inches(0.5),
       "Executive Summary", sz=26, bold=True, rgb=BLACK)
    rect(s2, Inches(0.6), Inches(1.35), Inches(2), Inches(0.06), GOLD)

    # ── Investment Thesis box ──
    # When the Recent Headlines sidebar has content, the thesis box shrinks
    # to make room for it on the right (~65/35 split). When there are no
    # headlines, the thesis takes the full 6.3" width — same as before.
    has_sidebar = bool(summary.headlines)
    thesis_w = Inches(4.1) if has_sidebar else Inches(6.3)

    tx(s2, Inches(0.6), Inches(1.6), Inches(2.5), Inches(0.3),
       "Investment Thesis", sz=12, bold=True, rgb=MUTED)
    rect(s2, Inches(0.6), Inches(1.95), thesis_w, Inches(3.2),
         THESIS_BG, THESIS_BORDER)
    rect(s2, Inches(0.6), Inches(1.95), Inches(0.06), Inches(3.2), GOLD)
    tx(s2, Inches(0.78), Inches(2.05), thesis_w - Inches(0.3), Inches(3.0),
       summary.thesis_text or "—", sz=11, rgb=BLACK, line_spacing=1.15)

    # ── Recent Headlines sidebar (slide 2 right column) ──
    # Editorial rule: news lives ONLY here, never stitched mid-thesis. The
    # thesis prose refers to themes only; specific headlines + dates +
    # sources live in this panel so reviewers can audit which headlines
    # informed the writeup. Section is fully suppressed when empty.
    if has_sidebar:
        sb_x = Inches(4.85)
        sb_y = Inches(1.95)
        sb_w = Inches(2.05)
        sb_h = Inches(3.2)
        rect(s2, sb_x, sb_y, sb_w, sb_h, WHITE, THESIS_BORDER)
        tx(s2, sb_x + Inches(0.12), sb_y + Inches(0.08), sb_w - Inches(0.2),
           Inches(0.22), "RECENT HEADLINES", sz=8, bold=True, rgb=MUTED)
        # Each headline gets ~0.7" — title (2 lines max) + date·source line.
        for i, h in enumerate(summary.headlines[:4]):
            row_y = sb_y + Inches(0.32 + i * 0.7)
            tx(s2, sb_x + Inches(0.12), row_y, sb_w - Inches(0.2),
               Inches(0.42), h.headline, sz=8, rgb=BLACK, line_spacing=1.05)
            meta_parts = []
            if h.date:
                meta_parts.append(h.date)
            if h.source:
                meta_parts.append(h.source)
            if meta_parts:
                tx(s2, sb_x + Inches(0.12), row_y + Inches(0.45),
                   sb_w - Inches(0.2), Inches(0.16),
                   " · ".join(meta_parts), sz=7, rgb=MUTED)

    # ── Charts ──
    chart_row_y = Inches(5.35)
    chart_row_h = Inches(2.15)
    try:
        if summary.income_chart and summary.income_chart.periods and (
            any(summary.income_chart.revenue) or any(summary.income_chart.net_income)
        ):
            tx(s2, Inches(0.6), Inches(5.0), Inches(3), Inches(0.3),
               "Income Statement Evolution", sz=11, bold=True, rgb=MUTED)
            ic = summary.income_chart
            build_revenue_ni_chart(
                s2,
                Inches(0.6), chart_row_y, Inches(3.2), chart_row_h,
                ic.periods, ic.revenue, ic.net_income,
                actuals_boundary=ic.actuals_boundary,
                # Currency is part of the units the chart axis already shows
                # via the "M" suffix; we don't double-print the ISO code.
                currency="",
                ebit_values=ic.ebit if ic.ebit else None,
            )

        if summary.pe_chart and summary.pe_chart.periods and any(
            v for v in summary.pe_chart.pe_values if v
        ):
            tx(s2, Inches(3.95), Inches(5.0), Inches(3), Inches(0.3),
               "P/E Multiple", sz=11, bold=True, rgb=MUTED)
            build_pe_chart(
                s2,
                Inches(3.95), chart_row_y, Inches(3.0), chart_row_h,
                summary.pe_chart.periods, summary.pe_chart.pe_values,
                five_yr_avg=summary.pe_chart.five_yr_avg,
            )

        # Source chip (top-right of chart row).
        chart_source = (
            (summary.income_chart.source_label if summary.income_chart else "")
            or "MarketScreener"
        )
        tx(s2, Inches(5.8), Inches(5.0), Inches(1.2), Inches(0.3),
           f"Source: {chart_source}", sz=7, rgb=MUTED, al=PP_ALIGN.RIGHT)
    except Exception as exc:
        # Never let chart failures break the rest of the slide. The earlier
        # contract said the same — keep going so the deck still ships.
        import logging
        logging.getLogger(__name__).warning("Chart rendering failed: %s", exc)

    # ── Key Expectations cards ──
    tx(s2, Inches(0.6), Inches(7.65), Inches(4), Inches(0.3),
       "Key Expectations", sz=14, bold=True, rgb=BLACK)
    cw = Inches(2.0)
    cg = Inches(0.15)
    for i, card in enumerate(summary.cards[:3]):
        x = Inches(0.6) + i * (cw + cg)
        rect(s2, x, Inches(8.0), cw, Inches(0.85), WHITE, THESIS_BORDER)
        # Label line — append the currency unit when present so the card
        # is self-describing (e.g. "Revenue · SARM").
        label_text = f"{card.label} · {card.unit}" if card.unit else card.label
        tx(s2, x + Inches(0.15), Inches(8.08), cw - Inches(0.3), Inches(0.2),
           label_text, sz=9, rgb=MUTED)
        tx(s2, x + Inches(0.15), Inches(8.3), cw - Inches(0.3), Inches(0.3),
           card.value_str, sz=18, bold=True, rgb=BLACK)
        if card.delta_str and card.delta_str != "—":
            tx(s2, x + Inches(0.15), Inches(8.62), cw - Inches(0.3), Inches(0.2),
               card.delta_str, sz=10, bold=True,
               rgb=_delta_color(card.delta_pct))

    # ── What to Watch ──
    # Editorial decision (this rewrite): no placeholder list. If the section
    # has no real content, the heading is suppressed entirely.
    if summary.what_to_watch:
        tx(s2, Inches(0.6), Inches(9.1), Inches(4), Inches(0.3),
           "What to Watch", sz=14, bold=True, rgb=BLACK)
        for i, item in enumerate(summary.what_to_watch[:4]):
            y = Inches(9.45) + Inches(i * 0.38)
            tx(s2, Inches(0.6), y, Inches(0.3), Inches(0.3),
               str(i + 1), sz=11, bold=True, rgb=GOLD)
            tx(s2, Inches(0.95), y, Inches(5.9), Inches(0.3),
               item, sz=11, rgb=BLACK)

    # ── Catalysts & Risks ──
    # Same editorial rule: omit when both lists are empty rather than ship
    # generic placeholders.
    if summary.catalysts or summary.risks:
        tx(s2, Inches(0.6), Inches(11.1), Inches(4), Inches(0.3),
           "Catalysts & Risks", sz=14, bold=True, rgb=BLACK)
        cbw = Inches(3.05)
        if summary.catalysts:
            rect(s2, Inches(0.6), Inches(11.45), cbw, Inches(0.9),
                 CATALYST_BG, THESIS_BORDER)
            rect(s2, Inches(0.6), Inches(11.45), Inches(0.06), Inches(0.9),
                 CATALYST_BAR)
            tx(s2, Inches(0.75), Inches(11.48), cbw - Inches(0.2),
               Inches(0.18), "CATALYSTS", sz=8, bold=True, rgb=CATALYST_BAR)
            tx(s2, Inches(0.75), Inches(11.65), cbw - Inches(0.2),
               Inches(0.65),
               "↑ " + "\n↑ ".join(summary.catalysts),
               sz=8, rgb=BLACK, word_wrap=False, line_spacing=0.9)
        if summary.risks:
            rx = Inches(3.85)
            rect(s2, rx, Inches(11.45), cbw, Inches(0.9),
                 RISK_BG, THESIS_BORDER)
            rect(s2, rx, Inches(11.45), Inches(0.06), Inches(0.9), RISK_BAR)
            tx(s2, rx + Inches(0.15), Inches(11.48), cbw - Inches(0.2),
               Inches(0.18), "KEY RISKS", sz=8, bold=True, rgb=RISK_BAR)
            tx(s2, rx + Inches(0.15), Inches(11.65), cbw - Inches(0.2),
               Inches(0.65),
               "↓ " + "\n↓ ".join(summary.risks),
               sz=8, rgb=BLACK, word_wrap=False, line_spacing=0.9)
