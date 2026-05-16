"""
Jabal Asset Management — Slide 1 (Snapshot) renderer.

Renders the first slide of the new 3-slide preview deck following the
spec in `docs/stage2/design_spec.md`. Reads exclusively from
`canonical_store` — never calls a provider directly. Returns the
pptx Slide object so a higher-level builder can stack slides 1/2/3.

Layout (top→bottom):
  1. Header strip
  2. Title block (kicker, company name, meta line, period subtitle)
  3. Analyst consensus row (3 cards)
  4. Key data row (6 metric blocks)
  5. Recent performance row (6 colored deltas)
  6. 52-week range bar
  7. Analyst highlights (5 pill rows)
  8. Footer
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from typing import Optional

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from src.services.jabal_design_tokens import (
    BLACK, GRAY, MUTED, GOLD, POS, NEG, CARD, WHITE,
    FONT_DISPLAY, FONT_UI,
    SZ_HERO, SZ_KICKER, SZ_VALUE, SZ_VALUE_LG, SZ_LABEL, SZ_BODY,
    SZ_META, SZ_HEADER, SZ_FOOTER, SZ_BULLET_PILL, SZ_TAB_NUM,
    PAGE_W_IN, PAGE_H_IN, MARGIN_L, MARGIN_R, CONTENT_W,
    RULE_THICK_PT, BORDER_THICK_PT, LEFT_ACCENT_W_IN,
    in_, signed_color,
)
from src.services.canonical_store import get_all_fields, CanonicalValue


# ── Low-level primitives ────────────────────────────────────

def _text(slide, left, top, width, height, text, *,
          font=FONT_UI, size=SZ_BODY, bold=False, color=BLACK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, all_caps=False,
          letter_spacing=None):
    """Insert a text box with one paragraph + one run. Returns the shape."""
    tb = slide.shapes.add_textbox(in_(left), in_(top), in_(width), in_(height))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text.upper() if all_caps else text
    r.font.name = font
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    if letter_spacing is not None:
        # python-pptx doesn't expose spc directly via property; set on xml
        from pptx.oxml.ns import qn
        r._r.get_or_add_rPr().set("spc", str(letter_spacing))
    return tb


def _hrule(slide, left, top, width, color=MUTED, thick_pt=RULE_THICK_PT):
    """Thin horizontal rule. Implemented as a 0.02"-tall rect with fill."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  in_(left), in_(top), in_(width), in_(0.005))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.shadow.inherit = False
    return shp


def _card(slide, left, top, width, height, *, fill=WHITE, border=MUTED,
          left_accent=GOLD, accent_w=LEFT_ACCENT_W_IN):
    """Bordered card with optional left-edge gold accent strip. Returns the
    main card shape (for layering text on top)."""
    # Main body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   in_(left), in_(top), in_(width), in_(height))
    body.fill.solid()
    body.fill.fore_color.rgb = fill
    body.line.color.rgb = border
    body.line.width = Pt(BORDER_THICK_PT)
    body.shadow.inherit = False
    # Left accent strip
    if left_accent is not None and accent_w > 0:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         in_(left), in_(top),
                                         in_(accent_w), in_(height))
        accent.fill.solid()
        accent.fill.fore_color.rgb = left_accent
        accent.line.fill.background()
        accent.shadow.inherit = False
    return body


def _metric_block(slide, left, top, width, label, value, *,
                  value_color=BLACK, value_size=SZ_VALUE):
    """Label-over-value metric primitive. label in muted 8.5pt caps,
    value in 14pt black."""
    _text(slide, left, top, width, 0.18, label, size=SZ_LABEL,
          color=MUTED, all_caps=True)
    _text(slide, left, top + 0.20, width, 0.32, value,
          size=value_size, color=value_color, bold=False)


def _section_label(slide, left, top, width, text):
    """10.5pt all-caps gray label with a horizontal rule above it."""
    _hrule(slide, left, top, width, color=MUTED, thick_pt=0.5)
    _text(slide, left, top + 0.10, width, 0.22, text,
          size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)


# ── Slide 1 sections ──────────────────────────────────────────

def _header_strip(slide, page_num: int, page_title: str):
    _text(slide, MARGIN_L, 0.32, 1.5, 0.28, "JABAL",
          font=FONT_DISPLAY, size=Pt(15), bold=True, color=BLACK)
    _text(slide, MARGIN_L, 0.58, 2.0, 0.18, "ASSET MANAGEMENT",
          size=Pt(8.5), color=GRAY, all_caps=True)
    _text(slide, 3.05, 0.36, 4.0, 0.22,
          f"PAGE {page_num}  ·  {page_title.upper()}",
          size=SZ_HEADER, color=BLACK, all_caps=True, align=PP_ALIGN.RIGHT)
    _text(slide, 3.05, 0.58, 4.0, 0.18, "INSTITUTIONAL RESEARCH  ·  EQUITY",
          size=Pt(8.5), color=MUTED, all_caps=True, align=PP_ALIGN.RIGHT)
    _hrule(slide, MARGIN_L, 0.88, CONTENT_W)


def _title_block(slide, company_name: str, meta_line: str, period: str):
    _text(slide, MARGIN_L, 1.08, CONTENT_W, 0.22, "EARNINGS PREVIEW NOTE",
          size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    _text(slide, MARGIN_L, 1.40, CONTENT_W, 0.62, company_name,
          font=FONT_DISPLAY, size=SZ_HERO, bold=False, color=BLACK)
    _text(slide, MARGIN_L, 2.08, CONTENT_W, 0.22, meta_line,
          size=SZ_META, color=GRAY)
    _text(slide, MARGIN_L, 2.36, CONTENT_W, 0.30, period,
          size=Pt(13), color=BLACK, bold=True)


def _consensus_row(slide, top: float, rating: str, n_analysts: int,
                    target_price: str, upside_pct: Optional[float]):
    _section_label(slide, MARGIN_L, top, CONTENT_W, "Analyst Consensus")
    row_top = top + 0.40
    card_w = (CONTENT_W - 0.22) / 3
    # Card 1: Rating
    _card(slide, MARGIN_L, row_top, card_w, 0.82, fill=CARD)
    _text(slide, MARGIN_L + 0.18, row_top + 0.08, card_w - 0.20, 0.18,
          f"RATING  ·  {n_analysts} ANALYSTS", size=SZ_LABEL, color=MUTED,
          all_caps=True)
    _text(slide, MARGIN_L + 0.18, row_top + 0.28, card_w - 0.20, 0.48,
          rating.upper(), size=SZ_VALUE_LG, color=BLACK, bold=True)
    # Card 2: Target Price
    c2_left = MARGIN_L + card_w + 0.11
    _card(slide, c2_left, row_top, card_w, 0.82, fill=CARD)
    _text(slide, c2_left + 0.18, row_top + 0.08, card_w - 0.20, 0.18,
          "TARGET PRICE", size=SZ_LABEL, color=MUTED, all_caps=True)
    _text(slide, c2_left + 0.18, row_top + 0.28, card_w - 0.20, 0.48,
          target_price, size=SZ_VALUE_LG, color=BLACK, bold=True)
    # Card 3: Upside
    c3_left = MARGIN_L + 2 * (card_w + 0.11)
    _card(slide, c3_left, row_top, card_w, 0.82, fill=CARD)
    _text(slide, c3_left + 0.18, row_top + 0.08, card_w - 0.20, 0.18,
          "UPSIDE TO TARGET", size=SZ_LABEL, color=MUTED, all_caps=True)
    up_str = "—" if upside_pct is None else f"{upside_pct:+.1f}%"
    _text(slide, c3_left + 0.18, row_top + 0.28, card_w - 0.20, 0.48,
          up_str, size=SZ_VALUE_LG, color=signed_color(upside_pct),
          bold=True)


def _key_data_row(slide, top: float, items: list[tuple[str, str]]):
    """Six (label, value) tuples spread across 6.6"."""
    _section_label(slide, MARGIN_L, top, CONTENT_W, "Key Data")
    row_top = top + 0.40
    n = len(items)
    col_w = CONTENT_W / n
    for i, (label, value) in enumerate(items):
        left = MARGIN_L + i * col_w
        _metric_block(slide, left, row_top, col_w - 0.10, label, value)


def _performance_row(slide, top: float, items: list[tuple[str, Optional[float]]]):
    _section_label(slide, MARGIN_L, top, CONTENT_W, "Recent Performance")
    row_top = top + 0.40
    n = len(items)
    col_w = CONTENT_W / n
    for i, (label, pct) in enumerate(items):
        left = MARGIN_L + i * col_w
        val_str = "—" if pct is None else f"{pct:+.1f}%"
        _metric_block(slide, left, row_top, col_w - 0.10, label, val_str,
                       value_color=signed_color(pct))


def _range_bar(slide, top: float, low: float, high: float, current: float,
                currency: str):
    _section_label(slide, MARGIN_L, top, CONTENT_W, "52-Week Range")
    bar_top = top + 0.48
    # Low/high labels at ends
    _text(slide, MARGIN_L, bar_top - 0.04, 0.9, 0.20,
          f"{currency} {low:,.2f}", size=Pt(10), color=GRAY)
    _text(slide, MARGIN_L + CONTENT_W - 0.9, bar_top - 0.04, 0.9, 0.20,
          f"{currency} {high:,.2f}", size=Pt(10), color=GRAY,
          align=PP_ALIGN.RIGHT)
    # Track
    track_left = MARGIN_L + 0.95
    track_w = CONTENT_W - 1.90
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    in_(track_left), in_(bar_top + 0.06),
                                    in_(track_w), in_(0.06))
    track.fill.solid()
    track.fill.fore_color.rgb = MUTED
    track.line.fill.background()
    # Fill from low to current
    if high > low:
        frac = max(0.0, min(1.0, (current - low) / (high - low)))
    else:
        frac = 0.5
    fill_w = track_w * frac
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   in_(track_left), in_(bar_top + 0.06),
                                   in_(max(fill_w, 0.02)), in_(0.06))
    fill.fill.solid()
    fill.fill.fore_color.rgb = GOLD
    fill.line.fill.background()
    # Diamond marker at current
    marker_size = 0.18
    marker = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                     in_(track_left + fill_w - marker_size / 2),
                                     in_(bar_top + 0.06 - marker_size / 2 + 0.03),
                                     in_(marker_size), in_(marker_size))
    marker.fill.solid()
    marker.fill.fore_color.rgb = BLACK
    marker.line.fill.background()
    # Current label below marker
    _text(slide, track_left + fill_w - 0.7, bar_top + 0.24, 1.4, 0.20,
          f"Current  {currency} {current:,.2f}",
          size=Pt(9), color=BLACK, align=PP_ALIGN.CENTER)


def _highlights_row(slide, top: float, items: list[tuple[str, str]]):
    """List of (category_pill, body) rows. Up to 5."""
    _section_label(slide, MARGIN_L, top, CONTENT_W,
                    "Analyst Highlights  ·  Key Points")
    row_top = top + 0.50
    for i, (cat, body) in enumerate(items[:5]):
        y = row_top + i * 0.42
        pill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       in_(MARGIN_L), in_(y),
                                       in_(0.85), in_(0.22))
        pill.fill.solid()
        pill.fill.fore_color.rgb = CARD
        pill.line.fill.background()
        _text(slide, MARGIN_L, y, 0.85, 0.22, cat,
              size=SZ_BULLET_PILL, color=GRAY, all_caps=True, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, MARGIN_L + 0.98, y - 0.02, CONTENT_W - 0.98, 0.32, body,
              size=SZ_BODY, color=BLACK)


def _footer(slide, page_num: int, total_pages: int, sources: str,
             analyst_name: str, gen_date: str):
    top = 12.47
    _hrule(slide, MARGIN_L, top, CONTENT_W)
    _text(slide, MARGIN_L, top + 0.07, CONTENT_W, 0.18,
          f"Source: {sources}  |  Generated {gen_date}  |  Analyst: {analyst_name}",
          size=SZ_FOOTER, color=GRAY)
    _text(slide, MARGIN_L, 12.90, 5.5, 0.18,
          "Jabal Asset Management  ·  Regulated by the Financial Services Authority of Oman",
          size=SZ_FOOTER, color=GRAY)
    _text(slide, 6.05, 12.90, 1.0, 0.18,
          f"{page_num} / {total_pages}",
          size=SZ_TAB_NUM, color=GRAY, align=PP_ALIGN.RIGHT)
    _text(slide, MARGIN_L, 13.08, CONTENT_W, 0.16,
          "CONFIDENTIAL  ·  For Institutional & Qualified Investors Only",
          size=Pt(7.5), color=MUTED, all_caps=True, align=PP_ALIGN.CENTER)


# ── Public entry point ────────────────────────────────────────

@dataclass
class SnapshotData:
    """The exact set of inputs Slide 1 needs. Built by the orchestrator
    from canonical_store + report metadata."""
    company_name: str
    ticker: str
    sector: str
    industry: str
    exchange: str
    period_label: str               # "Q2 2026 Earnings Preview"
    rating: str                     # "OUTPERFORM"
    n_analysts: int
    target_price_fmt: str           # "SAR 137.70"
    upside_pct: Optional[float]     # -4.5
    last_close_fmt: str
    market_cap_fmt: str
    report_date: str
    pe_fy_est_fmt: str
    div_yield_fmt: str
    currency: str
    perf_1d: Optional[float]
    perf_1w: Optional[float]
    perf_1m: Optional[float]
    perf_3m: Optional[float]
    perf_6m: Optional[float]
    perf_ytd: Optional[float]
    range_low: float
    range_high: float
    range_current: float
    highlights: list[tuple[str, str]]   # (CATEGORY, body) — max 5
    sources_line: str
    analyst_name: str
    gen_date: str
    total_pages: int = 3


def render_snapshot_slide(prs, data: SnapshotData):
    """Add slide 1 to an existing python-pptx Presentation. Returns the slide."""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    # Prefer a blank layout
    blank = next((L for L in prs.slide_layouts if L.name.lower() == "blank"),
                  layout)
    slide = prs.slides.add_slide(blank)

    # Sections
    _header_strip(slide, 1, "Snapshot")
    _title_block(
        slide, data.company_name,
        f"{data.ticker}  ·  {data.sector}  ·  {data.industry}  ·  {data.exchange}",
        data.period_label,
    )
    _consensus_row(slide, 2.96, data.rating, data.n_analysts,
                    data.target_price_fmt, data.upside_pct)
    _key_data_row(slide, 4.38, [
        ("LAST CLOSE", data.last_close_fmt),
        ("MARKET CAP", data.market_cap_fmt),
        ("REPORT DATE", data.report_date),
        ("P/E (FY EST)", data.pe_fy_est_fmt),
        ("DIV. YIELD", data.div_yield_fmt),
        ("CURRENCY", data.currency),
    ])
    _performance_row(slide, 5.42, [
        ("1 DAY", data.perf_1d),
        ("1 WEEK", data.perf_1w),
        ("1 MONTH", data.perf_1m),
        ("3 MONTHS", data.perf_3m),
        ("6 MONTHS", data.perf_6m),
        ("YTD", data.perf_ytd),
    ])
    _range_bar(slide, 6.46, data.range_low, data.range_high,
                data.range_current, data.currency)
    _highlights_row(slide, 7.62, data.highlights)
    _footer(slide, 1, data.total_pages, data.sources_line,
             data.analyst_name, data.gen_date)
    return slide


# ── Rating-label prettifier ───────────────────────────────────

def _pretty_rating(raw) -> str:
    """Normalise consensus-rating strings from any provider into the
    title-case form analysts use in print: e.g.
        'STRONG_BUY' -> 'Strong Buy'
        'OUTPERFORM' -> 'Outperform'
        'buy'        -> 'Buy'
        'hold/maintain' -> 'Hold/Maintain'
    Returns '' on empty/None input. Underscores and lower-case bleed
    through from Investing.com's enum (consensus_recommendation).
    """
    if raw is None:
        return ""
    s = str(raw).strip()
    if not s:
        return ""
    s = s.replace("_", " ").replace("-", " ")
    # Title-case but preserve runs of letters as words
    return " ".join(part.capitalize() for part in s.split())


# ── Sources-line builder ──────────────────────────────────────

def _sources_line(cv: dict) -> str:
    """Build the slide-footer "Source:" line.

    Lists every provider that contributed a value to ANY canonical field
    (i.e. the union of `sources_with_value`), not just the per-field
    winners. This means Investing.com / Yahoo / MarketScreener all show
    up when they fed the deck, even when only one of them won the
    reconciliation for a given cell.
    """
    seen: set[str] = set()
    for c in cv.values():
        winner = getattr(c, "canonical_source", "") or ""
        if winner:
            seen.add(winner)
        for s in (getattr(c, "sources_with_value", None) or []):
            if s:
                seen.add(s)
    # Prefer human-readable order: Yahoo, MarketScreener, Investing.com, others
    pretty = {
        "yahoo": "Yahoo Finance",
        "marketscreener": "MarketScreener",
        "investing": "Investing.com",
        "macro": "World Bank / IMF",
        "ishares": "iShares",
        "commodities": "World Bank / OPEC / EIA",
        "bloomberg": "Bloomberg",
        "ir_pdf": "Company IR",
    }
    ordered = sorted(seen, key=lambda s: (
        0 if s == "yahoo" else
        1 if s == "marketscreener" else
        2 if s == "investing" else
        3, s,
    ))
    labels = [pretty.get(s, s) for s in ordered]
    return ", ".join(labels) or "free-source stack"


# ── Highlight derivation ──────────────────────────────────────

def _derive_highlights(*, cv: dict, currency: str, current_price,
                         mcap, target_mean, upside_pct, pe_fwd,
                         div_yield, range_low, range_high,
                         n_analysts: int, rs: dict) -> list[tuple[str, str]]:
    """Derive 5 highlight pills from real numbers in canonical_store.
    Each row is anchored to a specific figure that also appears in the
    rest of the deck — keeps the slide internally consistent and avoids
    hardcoded boilerplate. Falls back to a short generic line only when
    no numeric anchor is available."""
    cur = currency or ""
    rows: list[tuple[str, str]] = []

    # EARNINGS — anchor on consensus rating + analyst count, with a fallback
    # to the Investing surprise track when present.
    rating_label = ""
    if isinstance(rs, dict):
        rating_label = _pretty_rating(rs.get("consensus"))
    if rating_label and n_analysts:
        rows.append(("EARNINGS", f"Consensus {rating_label} from {n_analysts} analysts covering."))
    elif n_analysts:
        rows.append(("EARNINGS", f"{n_analysts} analysts covering — view dispersion still narrow."))
    else:
        rows.append(("EARNINGS", "Awaiting next print; consensus build-up tracked across runs."))

    # VALUATION — P/E (FY est) is the most universally available figure.
    if isinstance(pe_fwd, (int, float)) and pe_fwd > 0:
        rows.append(("VALUATION", f"Forward P/E {float(pe_fwd):.1f}x — anchor for re-rate / de-rate debate."))
    elif isinstance(mcap, (int, float)) and mcap > 0:
        # No P/E available — fall back to market-cap framing.
        if mcap >= 1e12:
            rows.append(("VALUATION", f"Market cap {cur} {mcap/1e12:.2f}T — index-eligible scale."))
        elif mcap >= 1e9:
            rows.append(("VALUATION", f"Market cap {cur} {mcap/1e9:.1f}B — institutional-grade liquidity."))
        else:
            rows.append(("VALUATION", "Valuation context limited; awaiting MS / Investing refresh."))
    else:
        rows.append(("VALUATION", "Valuation context limited; awaiting MS / Investing refresh."))

    # POSITIONING — dividend yield is the strongest single signal for GCC banks,
    # SOEs, and dividend-heavy names. Fall back to current vs 52-week high
    # for growth names with no dividend.
    if isinstance(div_yield, (int, float)) and div_yield > 0:
        rows.append(("POSITIONING", f"Dividend yield {float(div_yield):.2f}% supports income mandate fit."))
    elif (isinstance(current_price, (int, float)) and current_price > 0
          and isinstance(range_high, (int, float)) and range_high > 0):
        gap = (current_price / range_high - 1.0) * 100
        rows.append(("POSITIONING", f"Trades {gap:+.1f}% versus 52-week high — entry-point context."))
    else:
        rows.append(("POSITIONING", "Range-trading context; refer to slide 3 for the 52-week band."))

    # WATCH — the target-vs-price gap is the cleanest forward-looking number
    # the audience asks about. Fall back to "next print awaited" line.
    if (isinstance(upside_pct, (int, float))
        and isinstance(target_mean, (int, float)) and target_mean > 0):
        rows.append(("WATCH", f"Target {cur} {target_mean:,.2f} ({upside_pct:+.1f}% vs last close)."))
    elif isinstance(target_mean, (int, float)) and target_mean > 0:
        rows.append(("WATCH", f"Consensus target sits at {cur} {target_mean:,.2f}."))
    else:
        rows.append(("WATCH", "Management commentary on forward outlook is the swing factor."))

    # RISK — analyst-distribution concentration is the most defensible
    # quantitative risk anchor (one-sided consensus = harder to surprise).
    # Some providers (notably MS for GCC names) publish only the consensus
    # label + analyst total without a buy/hold/sell breakdown; infer the
    # skew from the consensus label in that case so the pill still anchors
    # on a real signal.
    if isinstance(rs, dict):
        buy = int(rs.get("buy", 0) or 0)
        hold = int(rs.get("hold", 0) or 0)
        sell = int(rs.get("sell", 0) or 0)
        total = buy + hold + sell
        if total == 0 and n_analysts > 0:
            raw = (rs.get("consensus") or "").upper()
            pretty = _pretty_rating(rs.get("consensus"))
            if any(t in raw for t in ("BUY", "OUTPERFORM", "ACCUMULATE")):
                rows.append(("RISK",
                    f"Crowded long — consensus {pretty} across {n_analysts} analysts raises the expectations bar."))
            elif any(t in raw for t in ("SELL", "UNDERPERFORM", "REDUCE")):
                rows.append(("RISK",
                    f"Tape skewed bearish — consensus {pretty} across {n_analysts} analysts."))
            elif pretty:
                rows.append(("RISK",
                    f"View dispersion — consensus {pretty} across {n_analysts} analysts."))
            else:
                rows.append(("RISK", f"{n_analysts} analysts covering; breakdown not disclosed."))
        elif total > 0:
            denom = max(1, total)
            if sell == 0 and total >= 5:
                rows.append(("RISK", f"Sentiment one-sided — 0 sells across {total} analysts."))
            elif buy / denom >= 0.8 and total >= 5:
                rows.append(("RISK", f"Crowded long — {buy}/{total} buy ratings raise expectations bar."))
            elif sell >= 3:
                rows.append(("RISK", f"Tape skewed bearish — {sell}/{total} sell ratings."))
            else:
                rows.append(("RISK", f"Rating mix {buy}/{hold}/{sell} (buy/hold/sell) — view dispersion."))
        else:
            rows.append(("RISK", "Macro / sector sensitivity; refer to thesis on slide 2."))
    else:
        rows.append(("RISK", "Macro / sector sensitivity; refer to thesis on slide 2."))

    return rows


# ── Data adapter: canonical_store → SnapshotData ──────────────

def build_snapshot_data(ticker: str, *, analyst_name: str = "Jabal Research",
                          period_label: str = "Q2 2026 Earnings Preview",
                          report_date: str = "TBA",
                          highlights: Optional[list[tuple[str, str]]] = None,
                          ms_price_performance: Optional[dict] = None,
                          historical_override: Optional[dict] = None,
                          ) -> SnapshotData:
    """Translate canonical_store rows into the slide's input dataclass.
    Defensive against missing fields: every renderer-visible string has
    a sensible default ('—' for numerics, '' for text)."""
    cv = get_all_fields(ticker)

    def _val(field):
        c = cv.get(field)
        return c.value if c else None

    profile = _val("company_profile") or {}
    if not isinstance(profile, dict):
        profile = {}

    # Backfill profile from company_master when canonical_store hasn't been
    # populated by a recent yfinance probe (common on first-run tickers and
    # after a Render restart that wiped /tmp). This gives us at minimum a
    # sector / industry / exchange / country / company name to render in the
    # header subtitle — preventing the "· —" tail seen on fresh runs.
    try:
        from src.storage.db import load_company as _load_company
        cm = _load_company(ticker) or {}
        if cm:
            profile.setdefault("name", cm.get("company_name") or "")
            profile.setdefault("sector", cm.get("sector") or "")
            profile.setdefault("industry", cm.get("industry") or "")
            profile.setdefault("currency", cm.get("currency") or "")
            # Exchange suffix: friendly-name + country, e.g. "Tadawul (Saudi Arabia)".
            # The DB stores 3-letter codes (SAU, ADX, NSE, HKG, ...) — map to
            # human-readable bourse names. Unknown codes fall through to the raw code.
            _EX_NAMES = {
                "SAU": "Tadawul", "ADX": "ADX", "DFM": "DFM",
                "MSM": "MSX",  "DSM": "QSE", "BHB": "Bahrain Bourse",
                "KSE": "Boursa Kuwait", "EGX": "EGX",
                "NSE": "NSE",  "BSE": "BSE",
                "HKG": "HKEX", "SHA": "SSE", "SHZ": "SZSE",
                "TYO": "TSE",  "KRX": "KRX",
                "JNB": "JSE",
                "NMS": "NASDAQ", "NCM": "NASDAQ", "NGM": "NASDAQ",
                "NYQ": "NYSE", "NYS": "NYSE", "ASE": "NYSE American",
                "LON": "LSE", "PAR": "Euronext Paris",
                "AMS": "Euronext Amsterdam", "BRU": "Euronext Brussels",
                "FRA": "Frankfurt", "STO": "Nasdaq Stockholm",
                "ASX": "ASX", "TSE": "TSX",
            }
            xcode = (cm.get("exchange") or "").strip().upper()
            country = (cm.get("country") or "").strip()
            if xcode and country:
                profile["exchange"] = f"{_EX_NAMES.get(xcode, xcode)} ({country})"
            elif xcode:
                profile["exchange"] = _EX_NAMES.get(xcode, xcode)
    except Exception:
        pass
    last_price = _val("current_price")
    mcap = _val("market_cap")
    val_hist = _val("valuation_historical") or {}
    val_fwd = _val("valuation_forward") or {}
    target = _val("target_price")
    rating_split = _val("rating_split") or {}
    div_yield = _val("dividend_yield")
    hist_prices = _val("historical_prices") or {}
    # Prefer Investing-derived history when canonical_store is empty (the
    # GCC ex-Saudi tickers don't reach canonical_store via yfinance).
    if not (isinstance(hist_prices, dict) and hist_prices) and isinstance(historical_override, dict):
        hist_prices = historical_override

    # Currency: try profile first, else canonical_value units
    currency = (profile.get("currency") if isinstance(profile, dict) else None) or ""
    if not currency and cv.get("current_price"):
        currency = cv["current_price"].canonical_source[:3].upper()  # fallback

    # Format helpers
    def _money(x):
        if x is None:
            return "—"
        try:
            return f"{currency} {float(x):,.2f}" if currency else f"{float(x):,.2f}"
        except (TypeError, ValueError):
            return "—"

    # Different sources report market_cap in different units:
    #   MarketScreener: millions of local currency
    #   Yahoo / ADX / HKEX / NSE: raw local currency units
    mc_source = cv.get("market_cap").canonical_source if cv.get("market_cap") else ""
    mc_scale = 1_000_000.0 if mc_source == "marketscreener" else 1.0

    def _mc(x):
        if x is None:
            return "—"
        try:
            v = float(x) * mc_scale
        except (TypeError, ValueError):
            return "—"
        if v >= 1e12:
            return f"{currency} {v/1e12:.2f}T" if currency else f"{v/1e12:.2f}T"
        if v >= 1e9:
            return f"{currency} {v/1e9:.1f}B" if currency else f"{v/1e9:.1f}B"
        if v >= 1e6:
            return f"{currency} {v/1e6:.0f}M" if currency else f"{v/1e6:.0f}M"
        return f"{currency} {v:,.0f}"

    def _pct(x):
        if x is None:
            return "—"
        try:
            return f"{float(x):.2f}%"
        except (TypeError, ValueError):
            return "—"

    # Pull rating + analyst count
    rating = "—"
    n_analysts = 0
    if isinstance(rating_split, dict):
        # MS shape varies: ideally {"buy":6,"hold":3,"sell":1,"total":10,"consensus":"OUTPERFORM"}
        # but sometimes only {"consensus":"OUTPERFORM"}.
        n_analysts = int(rating_split.get("total", 0) or 0)
        rating = _pretty_rating(rating_split.get("consensus")) or "—"
    # Fall back to target_price.n_analysts if rating_split didn't carry the total.
    if not n_analysts and isinstance(target, dict):
        n_analysts = int(target.get("n_analysts", 0) or 0)

    # P/E forward — try MS forward dict first, then fall back to the
    # most recent valid value in valuation_historical.pe
    pe_fwd = None
    if isinstance(val_fwd, dict):
        for k in ("pe_fy1", "pe_2026", "pe_2027", "pe"):
            v = val_fwd.get(k)
            if isinstance(v, (int, float)):
                pe_fwd = v
                break
    elif isinstance(val_fwd, (int, float)):
        pe_fwd = val_fwd
    if pe_fwd is None and isinstance(val_hist, dict):
        for v in reversed(val_hist.get("pe", []) or []):
            if isinstance(v, (int, float)):
                pe_fwd = v
                break

    # Performance deltas — try Yahoo's hist_prices "perf_*" keys first;
    # fall back to MS price_performance block (perf_1d_pct etc.) when the
    # canonical historical_prices is empty (yfinance-blocked tickers).
    _ms_perf = (ms_price_performance or {}).get("performance") or {} \
        if isinstance(ms_price_performance, dict) else {}
    _ms_perf_keymap = {
        "perf_1d":  "perf_1d_pct",
        "perf_1w":  "perf_1w_pct",
        "perf_1m":  "perf_1m_pct",
        "perf_3m":  "perf_3m_pct",
        "perf_6m":  "perf_6m_pct",
        "perf_ytd": "perf_ytd_pct",
    }

    def _perf(key):
        if isinstance(hist_prices, dict):
            v = hist_prices.get(key)
            if v is not None:
                try: return float(v)
                except (TypeError, ValueError): pass
        ms_key = _ms_perf_keymap.get(key)
        if ms_key and isinstance(_ms_perf.get(ms_key), (int, float)):
            return float(_ms_perf[ms_key])
        return None

    # 52-week range
    low = high = current = None
    if isinstance(hist_prices, dict):
        low = hist_prices.get("range_52w_low")
        high = hist_prices.get("range_52w_high")
    try:
        current = float(last_price) if last_price is not None else None
    except (TypeError, ValueError):
        current = None
    if low is None or high is None or current is None:
        # Fallback: if we have only current, plot a degenerate range
        if current is not None:
            low = low if low is not None else current * 0.9
            high = high if high is not None else current * 1.1
        else:
            low, high, current = 0.0, 1.0, 0.5

    # MarketScreener's target_price comes as a dict
    # {"mean":29.86, "high":35.0, "low":26.5, "n_analysts":18}.
    # Yahoo / other providers may give a single float. Normalise.
    target_mean = None
    if isinstance(target, dict):
        target_mean = target.get("mean")
    elif isinstance(target, (int, float)):
        target_mean = target

    upside_pct = None
    try:
        if target_mean is not None and current is not None and current > 0:
            upside_pct = (float(target_mean) / current - 1.0) * 100
    except (TypeError, ValueError):
        pass

    target_fmt = "—"
    try:
        if target_mean is not None:
            target_fmt = (
                f"{currency} {float(target_mean):,.2f}" if currency
                else f"{float(target_mean):,.2f}"
            )
    except (TypeError, ValueError):
        pass

    # Compose meta-line pieces with safe fallbacks
    return SnapshotData(
        company_name=(profile.get("name") if isinstance(profile, dict) else None) or ticker,
        ticker=ticker,
        sector=(profile.get("sector") if isinstance(profile, dict) else None) or "—",
        industry=(profile.get("industry") if isinstance(profile, dict) else None) or "—",
        exchange=(profile.get("exchange") if isinstance(profile, dict) else None) or "—",
        period_label=period_label,
        rating=rating,
        n_analysts=n_analysts,
        target_price_fmt=target_fmt,
        upside_pct=upside_pct,
        last_close_fmt=_money(current),
        market_cap_fmt=_mc(mcap),
        report_date=report_date,
        pe_fy_est_fmt=("—" if pe_fwd is None else f"{float(pe_fwd):.1f}x"),
        div_yield_fmt=_pct(div_yield),
        currency=currency or "",
        perf_1d=_perf("perf_1d"),
        perf_1w=_perf("perf_1w"),
        perf_1m=_perf("perf_1m"),
        perf_3m=_perf("perf_3m"),
        perf_6m=_perf("perf_6m"),
        perf_ytd=_perf("perf_ytd"),
        range_low=float(low), range_high=float(high), range_current=float(current),
        highlights=highlights or _derive_highlights(
            cv=cv, currency=currency, current_price=current,
            mcap=mcap, target_mean=target_mean, upside_pct=upside_pct,
            pe_fwd=pe_fwd, div_yield=div_yield,
            range_low=low, range_high=high,
            n_analysts=n_analysts, rs=rating_split,
        ),
        sources_line=_sources_line(cv),
        analyst_name=analyst_name,
        gen_date=datetime.utcnow().strftime("%d %b %Y"),
    )
